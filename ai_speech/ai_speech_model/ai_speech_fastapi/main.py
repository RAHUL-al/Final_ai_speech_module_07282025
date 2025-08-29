from datetime import timezone
from google.cloud.firestore_v1.base_query import FieldFilter
from schemas import UserCreate, PasswordResetRequest, ForgotPasswordRequest, UserOut, UserUpdate, Token, LoginRequest, ForgotPasswordRequest, GeminiRequest, ChatRequest
from fastapi import Body
from fastapi.websockets import WebSocketState
import asyncio
from langchain_core.prompts import PromptTemplate
from langchain_ollama import ChatOllama
import aiohttp
from bs4 import BeautifulSoup
from langchain_core.messages import AIMessage, SystemMessage, HumanMessage
from langchain_core.output_parsers import StrOutputParser
from collections import Counter
import ujson as json
import redis
from datetime import datetime, timedelta
from fastapi import FastAPI, Depends, HTTPException, status, Request, WebSocket, WebSocketDisconnect, UploadFile, File, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from pydub import AudioSegment
from dotenv import load_dotenv
from urllib.parse import parse_qs, unquote
from concurrent.futures import ThreadPoolExecutor
from fastapi.responses import FileResponse, JSONResponse
import time
from pdf2image import convert_from_path
from typing import List, Optional, Any
from PIL import Image
import google.generativeai as genai
from pinecone import Pinecone
from auth import hash_password, verify_password, create_access_token
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
import torch
import logging
from urllib.parse import parse_qs,unquote
from jose import JWTError, jwt
from ai_speech_module import Topic, AdvancedAudioProcessor
from langchain_ollama import OllamaLLM
import re
from typing import List
from fastapi import FastAPI, UploadFile, File, Form
from pdf2image import convert_from_path
from PIL import Image
from docx2pdf import convert as docx_to_pdf
import os, uuid, shutil, logging, tempfile
from langchain.text_splitter import RecursiveCharacterTextSplitter
from fastapi import HTTPException, status
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
from PIL import Image, ImageDraw, ImageFont
import pandas as pd
import docx2txt
from pptx import Presentation
from firestore_models import FirestoreEssay
from firebase import db
from concurrent.futures import ProcessPoolExecutor
process_pool = ProcessPoolExecutor()
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks, WebSocket, WebSocketDisconnect
import aiohttp
import tempfile
from urllib.parse import parse_qs
from typing import List, Dict
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path
import docx2txt
from pptx import Presentation
import pandas as pd
from concurrent.futures import ThreadPoolExecutor, as_completed
import subprocess
from pdf2image import convert_from_path
from fastapi import HTTPException
from pathlib import Path
from PIL import Image
from typing import List, TypedDict, Annotated
import glob
from langchain.chains import LLMChain
import numpy as np
from scipy import signal
import librosa
from urllib.parse import parse_qs
from pydub.silence import detect_silence
from fastapi.concurrency import run_in_threadpool
from fastapi import BackgroundTasks
from fastapi_mail import FastMail, MessageSchema, ConnectionConfig

python311_path = "/home/ubuntu/ai_speech/venv/bin/python"
script_path = "/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi/layoutparser_file.py"


CPU_API_BASE = "http://13.200.201.10:8000"
model_name = OllamaLLM(model="mistral")
scraping_api_key = os.getenv("SCRAPINGDOG_API_KEY")

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_ENV = os.getenv("PINECONE_ENVIRONMENT")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
EMBEDDING_MODEL_NAME = "embaas/sentence-transformers-e5-large-v2"

SECRET_KEY = "jwt_secret_key"
ALGORITHM="HS256"


logging.basicConfig(
    filename='rag_log.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

load_dotenv()
app = FastAPI(title="FastAPI‑Firebase")

redis_client = redis.StrictRedis(
    host=os.getenv("REDIS_HOST"),
    port=int(os.getenv("REDIS_PORT")),
    username=os.getenv("REDIS_USERNAME"),
    password=os.getenv("REDIS_PASSWORD"),
    decode_responses=True
)

origins = ["https://llm.edusmartai.com","http://localhost:3000","http://localhost:5173"]

app.add_middleware(CORSMiddleware, allow_origins=origins, allow_credentials=True, allow_methods=["*"], allow_headers=["*"])


_ASCII_LETTER_RE = re.compile(r"[A-Za-z]")

_ARABIC_RANGES = [
    (0x0600, 0x06FF),
    (0x0750, 0x077F),
    (0x08A0, 0x08FF),
    (0xFB50, 0xFDFF),
    (0xFE70, 0xFEFF),
]

def contains_arabic_char(ch: str) -> bool:
    cp = ord(ch)
    for start, end in _ARABIC_RANGES:
        if start <= cp <= end:
            return True
    return False

def is_ascii_letter(ch: str) -> bool:
    return bool(_ASCII_LETTER_RE.match(ch))

def keep_token_if_english(token: str, letter_ratio_threshold: float = 0.6) -> bool:
    if token.strip() == "":
        return False

    if any(contains_arabic_char(ch) for ch in token):
        return False

    letters = [ch for ch in token if ch.isalpha()]
    if not letters:
        return False

    ascii_letters = sum(1 for ch in letters if is_ascii_letter(ch))
    ratio = ascii_letters / len(letters)
    return ratio >= letter_ratio_threshold

def filter_to_english(text: str, lower: bool = False) -> str:
    if lower:
        text = text.lower()

    tokens = re.split(r"\s+", text.strip())

    kept = []
    for tok in tokens:
        stripped = tok.strip("()[]{}\"'“”‘’.,;:?¡!—–")
        if keep_token_if_english(stripped):
            kept.append(stripped)

    return " ".join(kept).strip()

def get_user_from_redis_session(request: Request):
    token = request.headers.get("Authorization")
    if not token or not token.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid token")
    token = token.split(" ")[1]
    session_data = redis_client.get(f"session:{token}")
    if not session_data:
        raise HTTPException(status_code=401, detail="Session expired or invalid")
    return json.loads(session_data)

@app.post("/register", response_model=UserOut)
def register(user: UserCreate):
    user_ref = db.collection("users").where("username", "==", user.username).stream()
    if any(user_ref):
        raise HTTPException(400, "Username or email already exists")

    doc_ref = db.collection("users").add({
        "username": user.username,
        "email": user.email,
        "password": hash_password(user.password)
    })
    user_id = doc_ref[1].id
    return UserOut(id=user_id, username=user.username, email=user.email)



PASSWORD_RESET_TOKEN_EXPIRE_HOURS = 1

mail_conf = ConnectionConfig(
    MAIL_USERNAME=os.getenv("MAIL_USERNAME", "testampli2023@gmail.com"),
    MAIL_PASSWORD=os.getenv("MAIL_PASSWORD", "mulpeeeuolzidejx"),
    MAIL_FROM=os.getenv("MAIL_FROM", "testampli2023@gmail.com"),
    MAIL_PORT=int(os.getenv("MAIL_PORT", 587)),
    MAIL_SERVER=os.getenv("MAIL_SERVER", "smtp.gmail.com"),
    MAIL_STARTTLS=bool(os.getenv("MAIL_STARTTLS", True)),
    MAIL_SSL_TLS=bool(os.getenv("MAIL_SSL_TLS", False)),
    USE_CREDENTIALS=True,
    VALIDATE_CERTS=True
)


async def validate_reset_token(token: str):
    try:
        token_data = redis_client.get(f"session:{token}")

        logging.info(f"token data {token_data}")

        if not token_data:
            raise HTTPException(status_code=400, detail="Invalid or expired token")
        
        logging.info("token is validated -----------> done")
        return {"valid": True, "username": json.loads(token_data)["username"]}
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logging.error(f"Error validating token: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Invalid or expired token"
        )


def generate_reset_token(user_id: str):
    """Generate JWT token for password reset"""
    expires = datetime.now(timezone.utc) + timedelta(hours=PASSWORD_RESET_TOKEN_EXPIRE_HOURS)
    payload = {
        "sub": user_id,
        "exp": expires,
        "type": "password_reset"
    }
    return jwt.encode(payload, SECRET_KEY, algorithm=ALGORITHM)



async def send_reset_email(email: str, token: str, background_tasks: BackgroundTasks):
    reset_link = f"https://llm.edusmartai.com/reset-password?token={token}"
    message = MessageSchema(
        subject="Reset Your Password",
        recipients=[email],
        body=f"Click the link to reset your password: {reset_link}",
        subtype="plain"
    )
    fm = FastMail(mail_conf)
    background_tasks.add_task(fm.send_message, message)


@app.post("/forgot-password")
async def forgot_password(request: ForgotPasswordRequest, background_tasks: BackgroundTasks):
    try:
        if not request.email or "@" not in request.email:
            return {"detail": "If this email exists in our system, you'll receive a password reset link"}

        docs = db.collection("users").where(filter=FieldFilter("email", "==", request.email)).stream()
        user_doc = next(docs, None)

        if not user_doc:
            logging.info(f"Password reset requested for non-existent email: {request.email}")
            return {"detail": "If this email exists in our system, you'll receive a password reset link"}

        reset_token = generate_reset_token(user_doc.id)
        token_expiry = datetime.now(timezone.utc) + timedelta(hours=PASSWORD_RESET_TOKEN_EXPIRE_HOURS)

        redis_client.setex(
            f"reset_token:{reset_token}",
            int(timedelta(hours=PASSWORD_RESET_TOKEN_EXPIRE_HOURS).total_seconds()),
            json.dumps({
                "user_id": user_doc.id,
                "email": request.email,
                "expires_at": token_expiry.isoformat()
            })
        )

        await send_reset_email(request.email, reset_token, background_tasks)
        logging.info(f"Password reset token generated for {request.email}")

        return {"detail": "If this email exists in our system, you'll receive a password reset link"}

    except Exception as e:
        logging.error(f"Error in forgot-password for {request.email}: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="An error occurred while processing your request"
        )


@app.get("/validate-reset-token/{token}")
async def validate_reset_token_endpoint(token: str):
    try:
        token_data = redis_client.get(f"session:{token}")

        logging.info(f"token data {token_data}")

        if not token_data:
            raise HTTPException(status_code=400, detail="Invalid or expired token")
            
        return {"valid": True, "username": json.loads(token_data)["username"]}
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logging.error(f"Error validating token: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Invalid or expired token"
        )


@app.post("/reset-password")
async def reset_password(request: PasswordResetRequest = Body(...)):
    try:
        token_data_raw = redis_client.get(f"session:{request.token}")
        if not token_data_raw:
            raise HTTPException(status_code=400, detail="Invalid or expired token")

        logging.info(f"token data {token_data_raw}")

        token_data = json.loads(token_data_raw)
        user_id = token_data["user_id"]

        user_ref = db.collection("users").document(user_id)
        user_doc = user_ref.get()
        if not user_doc.exists:
            raise HTTPException(404, "User not found")

        user_ref.update({
            "password": hash_password(request.new_password),
            "updated_at": datetime.now(timezone.utc)
        })

        redis_client.delete(f"session:{request.token}")

        logging.info(f"Password reset successful for user {user_id}")
        return {"detail": "Password has been reset successfully"}

    except HTTPException as he:
        raise he
    except Exception as e:
        logging.error(f"Reset password error: {str(e)}", exc_info=True)
        raise HTTPException(500, "Failed to reset password")



@app.post("/login", response_model=Token)
def login(data: LoginRequest):
    docs = db.collection("users").where("username", "==", data.username).stream()
    user_doc = next(docs, None)
    if not user_doc or not verify_password(data.password, user_doc.to_dict()["password"]):
        raise HTTPException(status.HTTP_401_UNAUTHORIZED, detail="Bad credentials")

    token = create_access_token({"sub": user_doc.id})
    logging.info(f"token in login file : {token}")
    redis_client.setex(f"session:{token}", timedelta(hours=24), json.dumps({"user_id": user_doc.id, "username": data.username}))
    return Token(access_token=token, username=data.username)



@app.get("/logout")
def logout(request: Request):
    token = request.headers.get("Authorization")
    if token and token.startswith("Bearer "):
        redis_client.delete(f"session:{token.split(' ')[1]}")
    return {"detail": "Logged out"}

@app.get("/me", response_model=UserOut)
def me(user=Depends(get_user_from_redis_session)):
    doc = db.collection("users").document(user["user_id"]).get()
    if not doc.exists:
        raise HTTPException(404, "User not found")
    data = doc.to_dict()
    return UserOut(id=doc.id, username=data["username"], email=data["email"])



@app.post("/generate-prompt")
async def generate_prompt(data: GeminiRequest, user=Depends(get_user_from_redis_session)):
    
    url = f"https://en.wikipedia.org/wiki/{data.topic}"
    api_endpoint = f"https://api.scrapingdog.com/scrape?api_key={scraping_api_key}&url={url}"

    text = ""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(api_endpoint) as response:
                if response.status == 200:
                    html = await response.text()
                    soup = BeautifulSoup(html, "html.parser")
                    for script in soup(["script", "style", "noscript"]):
                        script.decompose()
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    logging.info(f"Error: {response.status} - {await response.text()}")
    except Exception as e:
        logging.exception(f"Failed to scrape data: {e}")

    prompt = (
        f"Generate a essay for a student in class {data.student_class} with a {data.accent} accent, "
        f"on the topic '{data.topic}', and the mood is '{data.mood}' and give me essay should be less than 400 words "
        f"and in response did not want \n\n or \n and also not want word count thanks you this type of stuff and used {text} "
        f"content for as updated data from internet and which is helpful in created essay and please give me output in paragraph format only not in points."
    )

    username = user.get("username")
    topic = Topic()
    response_text = await topic.topic_data_model_for_Qwen(username, prompt)

    essay_data = FirestoreEssay(
        username=username,
        user_id=user["user_id"],
        student_class=data.student_class,
        accent=data.accent,
        topic=data.topic,
        mood=data.mood,
        content=response_text
    )

    write_time, doc_ref = db.collection("essays").add(essay_data.to_dict())
    essay_id = doc_ref.id

    return JSONResponse(content={
        "response": response_text,
        "essay_id": essay_id
    })




@app.get("/overall-scoring-by-id")
async def overall_scoring_by_listening_module(essay_id: str):
    topic = Topic()
    try:
        result = await topic.overall_scoring_by_speech_module(essay_id)
        if result or result is None:

            return {"status": "success", "data": result}
        else:
            raise HTTPException(status_code=404, detail=f"No scoring found for essay_id: {essay_id}")
    except HTTPException:
        raise
    except Exception as e:
        last_exception = e

    raise HTTPException(status_code=500, detail=f"Internal error while fetching scoring: {str(last_exception)}")


@app.get("/overall-scoring-by-id-speech-module")
async def overall_scoring_by_listening_module(essay_id: str):
    topic = Topic()
    try:
        result = await topic.overall_scoring_by_speech_module(essay_id)
        if result or result is None:

            return {"status": "success", "data": result}
        else:
            raise HTTPException(status_code=404, detail=f"No scoring found for essay_id: {essay_id}")
    except HTTPException:
        raise
    except Exception as e:
        last_exception = e

    raise HTTPException(status_code=500, detail=f"Internal error while fetching scoring: {str(last_exception)}")
# @app.get("/overall-scoring-by-id-speech-module")
# async def overall_scoring_by_speech_module(essay_id: str):
#     topic = Topic()
#     total_attempts = 30
#     delay_seconds = 1

#     for attempt in range(total_attempts):
#         try:
#             result = await topic.overall_scoring_by_speech_module(essay_id)
#             if result is not None:
#                 return result
#         except Exception as e:
#             print(f"Attempt {attempt+1} failed: {e}")
#         await asyncio.sleep(delay_seconds)

#     raise HTTPException(status_code=404, detail=f"No scoring found for essay_id: {essay_id}")


# @app.get("/overall-scoring-by-id-speech-module")
# async def overall_scoring_by_speech_module(essay_id: str):
#     topic = Topic()
#     result = await topic.overall_scoring_by_speech_module(essay_id)
#     return result


async def final_listening__overall_scoring(result, transcribed_text,original_text,websocket,total_grammar_score):
    try:
        fluency = result.get("fluency")
        pronunciation = result.get("pronunciation")
        emotion = result.get("emotion")
        grammar_score = total_grammar_score

        prompt_template = PromptTemplate(template = """
            You are an expert English teacher providing comprehensive feedback on a student's listening comprehension exercise. Analyze the student's response in depth compared to the original essay.

            ORIGINAL ESSAY (Content to be comprehended):
            "{original_text}"

            STUDENT'S COMPLETE RESPONSE (What they actually spoke):
            "{transcribed_text}"

            PERFORMANCE METRICS:
            - Average Fluency: "{fluency}"
            - Average Pronunciation: "{pronunciation}" 
            - Dominant Emotion: "{emotion}"
            - Grammar Score: "{grammar_score}"

            CRITICAL ANALYSIS GUIDELINES:
            1. CONTENT COVERAGE ANALYSIS: 
            - Calculate approximate percentage of original content covered (word count, key ideas)
            - Identify if student captured main themes vs specific details
            - Note any content expansion or unnecessary additions

            2. SEMANTIC UNDERSTANDING (not just keyword matching):
            - Analyze if student understood concepts even with different vocabulary
            - Identify paraphrasing quality and conceptual accuracy
            - Detect misunderstandings or misinterpretations

            3. RESPONSE QUALITY ASSESSMENT:
            - Evaluate appropriateness of response length vs original content
            - Assess coherence and logical flow of ideas
            - Check if response addresses the core message

            4. SCORING ADJUSTMENT CONSIDERATIONS:
            - Adjust scores based on content coverage percentage
            - Consider difficulty level of the original text
            - Account for reasonable paraphrasing and synonym usage

            STRICT DETAILED FEEDBACK STRUCTURE JSON format always:

            {{
                "content_analysis": {{
                    "coverage_metrics": {{
                        "original_word_count": [calculate approximate word count],
                        "response_word_count": [calculate approximate word count],
                        "coverage_percentage": "X%",
                        "main_themes_captured": ["list main themes student understood"],
                        "major_omissions": ["important concepts student missed"]
                    }},
                    "understanding_depth": {{
                        "score": "X/10",
                        "explanation": "Detailed analysis of how well student grasped concepts beyond literal words. Include examples of good paraphrasing or conceptual understanding.",
                        "semantic_accuracy": "Assessment of whether meaning was preserved despite different wording",
                        "conceptual_gaps": "Any fundamental misunderstandings detected"
                    }},
                    "detail_retention": {{
                        "score": "X/10", 
                        "explanation": "Analysis of specific details, facts, or examples retained from original",
                        "key_details_captured": ["specific important details student mentioned"],
                        "important_details_missed": ["crucial details student omitted"]
                    }}
                }},
                
                "linguistic_performance": {{
                    "fluency_assessment": {{
                        "score": {fluency},
                        "detailed_analysis": "Specific examples of fluent segments and hesitations. Relate to content delivery.",
                        "pace_appropriateness": "Did speaking rate support comprehension of content?"
                    }},
                    "pronunciation_assessment": {{
                        "score": {pronunciation},
                        "strengths": ["Words pronounced correctly with examples"],
                        "improvement_areas": "Problematic words with phonetic corrections",
                        "impact_on_clarity": "How pronunciation affected content understanding"
                    }},
                    "grammar_assessment": {{
                        "score": {grammar_score},
                        "accuracy_analysis": "Specific grammatical errors with corrections",
                        "complexity_level": "Assessment of grammatical structures attempted",
                        "error_patterns": "Recurring grammar issues affecting meaning"
                    }},
                    "vocabulary_usage": {{
                        "appropriateness": "How well vocabulary matched the content level",
                        "synonym_usage": "Quality of word substitutions when paraphrasing",
                        "precision": "Use of specific vs generic terminology"
                    }}
                }},
                
                "technical_delivery": {{
                    "clarity_coherence": "How well organized was the response for listener understanding",
                    "emotion_appropriateness": "Did emotional tone match the content? {emotion}",
                    "comprehensibility_score": "Overall how understandable was the response (1-10)"
                }},
                
                "adaptive_scoring": {{
                    "content_adjusted_scores": {{
                        "understanding_score": "X/10 (adjusted for coverage and accuracy)",
                        "retention_score": "X/10 (adjusted for detail capture)",
                        "overall_comprehension": "X/10 (composite score)"
                    }},
                    "scoring_rationale": "Explanation of how coverage percentage affected final scores"
                }},
                
                "improvement_strategy": {{
                    "immediate_priority": "Most critical area needing improvement with specific examples",
                    "content_handling_tips": [
                        "How to better capture main ideas",
                        "Techniques for detail retention", 
                        "Paraphrasing strategies"
                    ],
                    "practice_recommendations": [
                        "Specific exercises for content comprehension",
                        "Listening retention drills",
                        "Summarization practice"
                    ]
                }},
                                         
                "overall_scores": {{
                    "fluency": "{fluency}",
                    "pronunciation": "{pronunciation}",
                    "grammar": "{grammar_score}",
                    "emotion": f"{emotion}"
                }},
                
                "strengths_highlight": [
                    "What student did well despite content coverage limitations",
                    "Effective paraphrasing examples",
                    "Good conceptual understanding demonstrations"
                ],
                
                "encouragement_message": "Motivational note acknowledging effort and suggesting growth path based on actual performance"
            }}

            SPECIFIC ANALYSIS REQUIREMENTS:
            1. For short responses: Focus on quality of what was said rather than quantity missed
            2. For keyword differences: Analyze if meaning was preserved with different terminology  
            3. For partial coverage: Assess depth of understanding in covered portions
            4. Always provide concrete examples from both original and student response
            5. Be constructive - suggest how to improve rather than just criticizing
            6. Highlight specific strengths to encourage student
            7. All things need to tell with full explanation not in one-two sentence

            Provide detailed, specific feedback that genuinely helps the student improve their listening comprehension skills.
            """,
            input_variables=["original_text", "transcribed_text", "fluency", "pronunciation", "emotion", "grammar_score","fluency","pronunciation","grammar_score","emotion"]
            )
        parser = StrOutputParser()
        chain = prompt_template | model | parser
        try:
            feedback = await chain.ainvoke({
                "original_text": original_text,
                "transcribed_text": transcribed_text,
                "fluency": fluency,
                "pronunciation": pronunciation,
                "emotion": emotion,
                "grammar_score": grammar_score,
            })

            if feedback:
                json_match = re.search(r'\{.*\}', feedback, re.DOTALL)
                if json_match:
                    try:
                        feedback_json = json.loads(json_match.group())
                        return feedback_json
                    except json.JSONDecodeError as je:
                        logging.error(f"JSON parsing failed: {je}\nRaw feedback: {feedback}")
                        return feedback
                else:
                    logging.error(f"No JSON object found in feedback\nRaw feedback: {feedback}")
                    return {"error": "No JSON object found in feedback", "raw_feedback": feedback}
            else:
                return {"error": "Empty feedback received"}

        except Exception as e:
            logging.error(f"Session finalization error: {e}")
            return {"error": str(e)}

    except Exception as e:
        logging.info(f"Session finalization error: {e}")




TEMP_DIR = os.path.abspath("audio_folder")
os.makedirs(TEMP_DIR, exist_ok=True)

@app.websocket("/ws/audio")
async def audio_ws(websocket: WebSocket):
    await websocket.accept()
    query_params = parse_qs(websocket.url.query)
    username = query_params.get("username", [None])[0]
    token = query_params.get("token", [None])[0]

    if not username or not token:
        await websocket.close(code=4001)
        logging.info("Username or token missing.")
        return

    logging.info(f"[WS] Authenticated connection from {username}")
    chunk_index = 0
    chunk_results = []
    text_output = []

    session_state = {
        'audio_buffer': AudioSegment.empty(),
        'silence_duration': 0.0,
        'is_speaking': False
    }

    silence_threshold = 1
    min_audio_length = 1.0

    date_str = datetime.now().strftime("%Y-%m-%d")
    user_dir = os.path.join(TEMP_DIR, username, date_str)
    os.makedirs(user_dir, exist_ok=True)

    final_output = os.path.join(user_dir, f"{username}_output.wav")
    transcript_path = os.path.join(user_dir, f"{username}_transcript.txt")

    if os.path.exists(final_output):
        os.remove(final_output)
    if os.path.exists(transcript_path):
        os.remove(transcript_path)

    total_grammar_score = 0
    grammar_score_count = 0
    average_scores = {"fluency": 0, "pronunciation": 0, "emotion": "neutral"}

    topic = Topic()
    aiohttp_session = aiohttp.ClientSession()

    try:
        while True:
            try:
                # ✅ Correct: explicitly receive raw audio bytes
                audio_bytes = await websocket.receive_bytes()
            except WebSocketDisconnect:
                logging.info(f"[WS] {username} disconnected gracefully.")
                break
            except Exception as e:
                logging.error(f"[WS] Error receiving audio: {e}")
                break

            audio = AudioSegment(
                data=audio_bytes,
                sample_width=2,
                frame_rate=16000,
                channels=1
            )
            session_state['audio_buffer'] += audio

            # detect silence
            silence_ranges = detect_silence(audio, min_silence_len=100, silence_thresh=-40)
            current_silence = sum((end - start) / 1000.0 for start, end in silence_ranges) if silence_ranges else 0.0

            try:
                temp_chunk_path = os.path.join(user_dir, f"temp_chunk_{chunk_index}.wav")
                audio.export(temp_chunk_path, format="wav")

                silvero = await topic.silvero_vad(temp_chunk_path)
                current_silence = silvero.get("silence_duration", current_silence)
                is_speaking = silvero.get("is_speaking", False)
                os.remove(temp_chunk_path)
            except Exception as e:
                logging.warning(f"Silvero VAD failed: {e}, fallback mode")
                is_speaking = current_silence < (len(audio) / 1000.0) * 0.7

            if is_speaking:
                session_state['silence_duration'] = 0.0
                session_state['is_speaking'] = True
            else:
                session_state['silence_duration'] += current_silence
                session_state['is_speaking'] = False

            # flush when silence threshold crossed
            if session_state['silence_duration'] >= silence_threshold and len(session_state['audio_buffer']) > min_audio_length * 1000:
                chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
                session_state['audio_buffer'].export(chunk_filename, format="wav")

                try:
                    transcribed_text = await topic.speech_to_text(chunk_filename, username)
                    output = filter_to_english(transcribed_text)
                    if not output:
                        continue

                    await websocket.send_json({"type": "transcription", "data": output})
                    logging.info(f"[WS] {username} transcribed: {output}")

                    grammar_score = await topic.grammar_checking(transcribed_text)
                    if grammar_score:
                        total_grammar_score += float(grammar_score)
                        grammar_score_count += 1

                    # call scoring APIs
                    fluency, pronunciation, emotion = 0, 0, "unknown"

                    async with aiohttp_session.post(
                        f"{CPU_API_BASE}/fluency-score",
                        json={"text": output}
                    ) as res:
                        fluency_data = await res.json()
                        fluency = fluency_data.get("fluency", 0)

                    # Pronunciation scoring
                    with open(chunk_filename, "rb") as f:
                        form = aiohttp.FormData()
                        form.add_field("file", f, filename=os.path.basename(chunk_filename), 
                                        content_type="audio/wav")
                        async with aiohttp_session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as res:
                            pron_data = await res.json()
                            pronunciation = pron_data.get("pronunciation", 0)

                    # async with aiohttp_session.post(
                    #     f"{CPU_API_BASE}/pronunciation-score",
                    #     data={"text": output}
                    # ) as res:
                    #     pron_data = await res.json()
                    #     pronunciation = pron_data.get("pronunciation", 0)

                    with open(chunk_filename, "rb") as f:
                        form = aiohttp.FormData()
                        form.add_field("file", f, filename=os.path.basename(chunk_filename), content_type="audio/wav")
                        async with aiohttp_session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as res:
                            emotion_data = await res.json()
                            emotion = emotion_data.get("emotion", "unknown")

                    topic.update_realtime_stats(fluency, pronunciation, emotion)
                    text_output.append(output)

                    chunk_result = {
                        "chunk_index": chunk_index,
                        "text": output,
                        "emotion": emotion,
                        "fluency": fluency,
                        "pronunciation": pronunciation,
                        "grammar_score": grammar_score,
                        "silence_duration": session_state['silence_duration'],
                        "file_path": chunk_filename
                    }
                    chunk_results.append(chunk_result)

                    # reset buffer
                    session_state['audio_buffer'] = AudioSegment.empty()
                    session_state['silence_duration'] = 0.0
                    session_state['is_speaking'] = False
                    chunk_index += 1

                except Exception as e:
                    logging.error(f"Error processing chunk {chunk_index}: {e}")

    finally:
        await aiohttp_session.close()

        # Only update DB if we got results
        if chunk_results:
            loop = asyncio.get_event_loop()
            await loop.run_in_executor(None, merge_chunks, chunk_results, final_output)

            with open(transcript_path, "w", encoding="utf-8") as f:
                f.write(" ".join(text_output).strip())

            try:
                essays_ref = db.collection("essays").where("username", "==", username)
                essays = essays_ref.stream()
                today = datetime.now().date()

                essay_list = []
                for essay in essays:
                    essay_data = essay.to_dict()
                    create_time = essay_data.get("timestamp", essay.create_time)
                    create_date = create_time.date() if hasattr(create_time, 'date') else create_time.to_python().date()
                    if create_date == today:
                        essay_list.append((essay.id, essay_data, create_time))

                if essay_list:
                    latest_essay_id, latest_essay_data, latest_create_time = max(essay_list, key=lambda x: x[2])
                    essay_ref = db.collection("essays").document(latest_essay_id)

                    average_scores = topic.get_average_realtime_scores()
                    if grammar_score_count > 0:
                        total_grammar_score = round(total_grammar_score / grammar_score_count, 2)

                    text_output_str = " ".join(text_output).strip()
                    original_text = latest_essay_data.get("content", "")

                    feedback_result = await final_listening__overall_scoring(
                        average_scores, text_output_str, original_text, websocket, total_grammar_score
                    )

                    essay_ref.update({
                        "chunks": chunk_results,
                        "average_scores": average_scores,
                        "overall_scoring": feedback_result
                    })
                    logging.info(f"Updated essay {latest_essay_id}")

            except Exception as e:
                logging.error(f"[Firestore Update Error] {e}")


def merge_chunks(chunk_files, final_output):
    logging.info("[Merge] Merging audio chunks...")
    combined = AudioSegment.empty()

    for chunk in chunk_files:
        file_path = chunk.get("file_path")
        if file_path and os.path.exists(file_path):
            try:
                audio = AudioSegment.from_file(file_path, format="wav")
                combined += audio
            except Exception as e:
                logging.warning(f"[Merge] Error reading {file_path}: {e}")
        else:
            logging.warning(f"[Merge] Skipping missing or invalid file: {file_path}")

    if len(combined) > 0:
        combined.export(final_output, format="wav")
        logging.info("[Merge] Merged audio file saved.")
    else:
        logging.warning("[Merge] No audio to merge.")

# @app.websocket("/ws/audio")
# async def audio_ws(websocket: WebSocket):
#     await websocket.accept()
#     query_params = parse_qs(websocket.url.query)
#     username = query_params.get("username", [None])[0]
#     token = query_params.get("token", [None])[0]

#     if not username or not token:
#         await websocket.close(code=4001)
#         logging.info("Username or token missing.")
#         return

#     logging.info(f"[WS] Authenticated connection from {username}")
#     chunk_index = 0
#     chunk_results = []
#     text_output = []
    
#     session_state = {
#         'audio_buffer': AudioSegment.empty(),
#         'silence_duration': 0.0,
#         'is_speaking': False
#     }

#     silence_threshold = 1
#     min_audio_length = 1.0

#     date_str = datetime.now().strftime("%Y-%m-%d")
#     user_dir = os.path.join(TEMP_DIR, username, date_str)
#     os.makedirs(user_dir, exist_ok=True)

#     final_output = os.path.join(user_dir, f"{username}_output.wav")
#     transcript_path = os.path.join(user_dir, f"{username}_transcript.txt")

#     if os.path.exists(final_output):
#         os.remove(final_output)
#     if os.path.exists(transcript_path):
#         os.remove(transcript_path)
#     total_grammar_score = 0
#     grammar_score_count = 0
#     average_scores = {"fluency": 0, "pronunciation": 0, "emotion": "neutral"}

#     try:
#         topic = Topic()
#         while True:
#             message = await websocket.receive()

#             if message["type"] == "websocket.disconnect":
#                 print(f"[WS] {username} disconnected.")
#                 break

#             if message["type"] == "websocket.receive" and "bytes" in message:
#                 audio = AudioSegment(
#                     data=message["bytes"],
#                     sample_width=2,
#                     frame_rate=16000,
#                     channels=1
#                 )
                
#                 session_state['audio_buffer'] += audio
                
#                 silence_ranges = detect_silence(
#                     audio, 
#                     min_silence_len=100,
#                     silence_thresh=-40
#                 )
                
#                 current_silence = 0.0
#                 if silence_ranges:
#                     current_silence = sum((end - start) / 1000.0 for start, end in silence_ranges)
                
#                 try:
#                     temp_chunk_path = os.path.join(user_dir, f"temp_chunk_{chunk_index}.wav")
#                     audio.export(temp_chunk_path, format="wav")
                    
#                     silvero = await topic.silvero_vad(temp_chunk_path)
#                     current_silence = silvero.get("silence_duration", current_silence)
#                     is_speaking = silvero.get("is_speaking", False)
                    
#                     os.remove(temp_chunk_path)
                    
#                 except Exception as e:
#                     logging.warning(f"Silvero VAD failed: {e}, using fallback silence detection")
#                     is_speaking = current_silence < (len(audio) / 1000.0) * 0.7
                
#                 if is_speaking:
#                     session_state['silence_duration'] = 0.0
#                     session_state['is_speaking'] = True
#                 else:
#                     session_state['silence_duration'] += current_silence
#                     session_state['is_speaking'] = False
                
#                 if (session_state['silence_duration'] >= silence_threshold and 
#                     len(session_state['audio_buffer']) > min_audio_length * 1000):
                    
#                     chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
#                     session_state['audio_buffer'].export(chunk_filename, format="wav")
                    
#                     try:
#                         async with aiohttp.ClientSession() as session:
#                             transcribed_text = await topic.speech_to_text(chunk_filename, username)
                            
#                             output = filter_to_english(transcribed_text)
#                             if not output:
#                                 return
#                             await websocket.send_json({"type": "transcription", "data": output})
#                             logging.info(f"[WS] {username} transcribed text: {output}")
#                             grammar_score = await topic.grammar_checking(transcribed_text)
#                             if grammar_score:
#                                 total_grammar_score += float(grammar_score)
#                                 grammar_score_count += 1

#                             with open(chunk_filename, "rb") as f:
#                                 form = aiohttp.FormData()
#                                 form.add_field("file", f, filename=os.path.basename(chunk_filename), 
#                                              content_type="audio/wav")
#                                 async with session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as res:
#                                     emotion_data = await res.json()
#                                     emotion = emotion_data.get("emotion", "unknown")

#                             async with session.post(
#                                 f"{CPU_API_BASE}/fluency-score",
#                                 json={"text": output}
#                             ) as res:   
#                                 fluency_data = await res.json()
#                                 fluency = fluency_data.get("fluency", 0)

#                             with open(chunk_filename, "rb") as f:
#                                 form = aiohttp.FormData()
#                                 form.add_field("file", f, filename=os.path.basename(chunk_filename), 
#                                              content_type="audio/wav")
#                                 async with session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as res:
#                                     pron_data = await res.json()
#                                     pronunciation = pron_data.get("pronunciation", 0)

#                         topic.update_realtime_stats(fluency, pronunciation, emotion)
#                         text_output.append(output)

#                         chunk_result = {
#                             "chunk_index": chunk_index,
#                             "text": output,
#                             "emotion": emotion,
#                             "fluency": fluency,
#                             "pronunciation": pronunciation,
#                             "silence_duration": session_state['silence_duration'],
#                             "file_path": chunk_filename
#                         }

#                         logging.info(f"[Chunk {chunk_index}] {chunk_result}")
#                         chunk_results.append(chunk_result)
                        
#                         session_state['audio_buffer'] = AudioSegment.empty()
#                         session_state['silence_duration'] = 0.0
#                         session_state['is_speaking'] = False
#                         chunk_index += 1
                        
#                     except Exception as e:
#                         logging.error(f"Error processing chunk {chunk_index}: {e}")
                
#     except WebSocketDisconnect:
#         logging.warning(f"[WS] {username} forcibly disconnected.")
#     except Exception as e:
#         logging.error(f"[WS] Unexpected error for {username}: {e}")

#     finally:
#         # Process any remaining audio in buffer
#         if len(session_state['audio_buffer']) > min_audio_length * 1000:
#             chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
#             session_state['audio_buffer'].export(chunk_filename, format="wav")
            
#             try:
#                 async with aiohttp.ClientSession() as session:
#                     transcribed_text = await topic.speech_to_text(chunk_filename, username)
#                     output = filter_to_english(transcribed_text)
                    
#                     if output:
#                         await websocket.send_json({"type": "transcription", "data": output})
#                         text_output.append(output)
                        
#                         grammar_score = await topic.grammar_checking(transcribed_text)
#                         if grammar_score:
#                             total_grammar_score += float(grammar_score)
#                             grammar_score_count += 1
                            
#                         chunk_result = {
#                             "chunk_index": chunk_index,
#                             "text": output,
#                             "emotion": "unknown",
#                             "fluency": 0,
#                             "pronunciation": 0,
#                             "grammar_score": grammar_score,
#                             "silence_duration": 0,
#                             "file_path": chunk_filename
#                         }
#                         chunk_results.append(chunk_result)
#             except Exception as e:
#                 logging.error(f"Error processing final chunk: {e}")
        
#         loop = asyncio.get_event_loop()
#         await loop.run_in_executor(None, merge_chunks, chunk_results, final_output)

#         with open(transcript_path, "w", encoding="utf-8") as f:
#             f.write(" ".join(text_output).strip())

#         try:
#             essays_ref = db.collection("essays").where("username", "==", username)
#             essays = essays_ref.stream()
#             today = datetime.now().date()
            
#             essay_list = []
#             for essay in essays:
#                 essay_data = essay.to_dict()
                
#                 # Get the creation timestamp from the document
#                 create_time = essay_data.get("timestamp")
#                 if not create_time:
#                     # If no timestamp field, use the document creation time
#                     create_time = essay.create_time
                
#                 # Convert to date for comparison
#                 if hasattr(create_time, 'date'):
#                     create_date = create_time.date()
#                 else:
#                     # Handle Firestore Timestamp objects
#                     create_date = create_time.to_python().date()
                
#                 if create_date == today:
#                     essay_list.append((essay.id, essay_data, create_time))
            
#             if essay_list:
#                 # Get the latest essay by creation time
#                 latest_essay_id, latest_essay_data, latest_create_time = max(
#                     essay_list, key=lambda x: x[2]
#                 )
                
#                 # Get a reference to the latest essay document
#                 essay_ref = db.collection("essays").document(latest_essay_id)

#                 average_scores = topic.get_average_realtime_scores()
#                 logging.info(f"Average Scores: {average_scores}")

#                 text_output_str = " ".join(text_output).strip()

#                 # Get the original text from the essay data
#                 original_text = latest_essay_data.get("content", "")
                
#                 # Calculate average grammar score
#                 if grammar_score_count > 0:
#                     total_grammar_score = round(total_grammar_score / grammar_score_count, 2)
#                 else:
#                     total_grammar_score = 0

#                 feedback_result = await final_listening__overall_scoring(
#                     average_scores, 
#                     text_output_str, 
#                     original_text,
#                     websocket,
#                     total_grammar_score
#                 )

#                 if "```json" in feedback_result:
#                     feedback_result = feedback_result.replace("```json", "")
#                 if "```" in feedback_result:
#                     feedback_result = feedback_result.replace("```", "")

#                 logging.info(f"Final Feedback: {feedback_result}")
                
#                 essay_ref.update({
#                     "chunks": chunk_results,
#                     "average_scores": average_scores,
#                     "overall_scoring": feedback_result
#                 })
#                 logging.info(f"Updated essay {latest_essay_id}")

#         except Exception as e:
#             logging.error(f"[Firestore Update Error] {e}")

#         finally:
#             # Clean up chunk files
#             for file in os.listdir(user_dir):
#                 if file.startswith("chunk_") and file.endswith(".wav"):
#                     try:
#                         os.remove(os.path.join(user_dir, file))
#                     except Exception as e:
#                         logging.warning(f"Failed to remove {file}: {e}")




# def merge_chunks(chunk_files, final_output):
#     logging.info("[Merge] Merging audio chunks...")
#     combined = AudioSegment.empty()

#     for chunk in chunk_files:
#         file_path = chunk.get("file_path")
#         if file_path and os.path.exists(file_path):
#             try:
#                 audio = AudioSegment.from_file(file_path, format="wav")
#                 combined += audio
#             except Exception as e:
#                 logging.warning(f"[Merge] Error reading {file_path}: {e}")
#         else:
#             logging.warning(f"[Merge] Skipping missing or invalid file: {file_path}")

#     combined.export(final_output, format="wav")
#     logging.info("[Merge] Merged audio file saved.")
















# @app.websocket("/ws/audio")
# async def audio_ws(websocket: WebSocket):
#     await websocket.accept()
#     query_params = parse_qs(websocket.url.query)
#     username = query_params.get("username", [None])[0]
#     token = query_params.get("token", [None])[0]

#     if not username or not token:
#         await websocket.close(code=4001)
#         logging.info("Username or token missing.")
#         return

#     logging.info(f"[WS] Authenticated connection from {username}")
#     chunk_index = 0
#     chunk_results = []
#     text_output = []
    
#     session_state = {
#         'audio_buffer': AudioSegment.empty(),
#         'silence_duration': 0.0,
#         'is_speaking': False
#     }

#     silence_threshold = 2
#     min_audio_length = 1.0

#     date_str = datetime.now().strftime("%Y-%m-%d")
#     user_dir = os.path.join(TEMP_DIR, username, date_str)
#     os.makedirs(user_dir, exist_ok=True)

#     final_output = os.path.join(user_dir, f"{username}_output.wav")
#     transcript_path = os.path.join(user_dir, f"{username}_transcript.txt")

#     if os.path.exists(final_output):
#         os.remove(final_output)
#     if os.path.exists(transcript_path):
#         os.remove(transcript_path)
    
#     # Initialize variables with default values
#     total_grammar_score = 0
#     grammar_score_count = 0
#     average_scores = {"fluency": 0, "pronunciation": 0, "emotion": "neutral"}
#     average_grammar_score = 0

#     try:
#         topic = Topic()
#         while True:
#             message = await websocket.receive()

#             if message["type"] == "websocket.disconnect":
#                 print(f"[WS] {username} disconnected.")
#                 break

#             if message["type"] == "websocket.receive" and "bytes" in message:
#                 audio = AudioSegment(
#                     data=message["bytes"],
#                     sample_width=2,
#                     frame_rate=16000,
#                     channels=1
#                 )
                
#                 session_state['audio_buffer'] += audio
                
#                 silence_ranges = detect_silence(
#                     audio, 
#                     min_silence_len=100,
#                     silence_thresh=-40
#                 )
                
#                 current_silence = 0.0
#                 if silence_ranges:
#                     current_silence = sum((end - start) / 1000.0 for start, end in silence_ranges)
                
#                 try:
#                     temp_chunk_path = os.path.join(user_dir, f"temp_chunk_{chunk_index}.wav")
#                     audio.export(temp_chunk_path, format="wav")
                    
#                     silvero = await topic.silvero_vad(temp_chunk_path)
#                     current_silence = silvero.get("silence_duration", current_silence)
#                     is_speaking = silvero.get("is_speaking", False)
                    
#                     os.remove(temp_chunk_path)
                    
#                 except Exception as e:
#                     logging.warning(f"Silvero VAD failed: {e}, using fallback silence detection")
#                     is_speaking = current_silence < (len(audio) / 1000.0) * 0.7
                
#                 if is_speaking:
#                     session_state['silence_duration'] = 0.0
#                     session_state['is_speaking'] = True
#                 else:
#                     session_state['silence_duration'] += current_silence
#                     session_state['is_speaking'] = False
                
#                 if (session_state['silence_duration'] >= silence_threshold and 
#                     len(session_state['audio_buffer']) > min_audio_length * 1000):
                    
#                     chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
#                     session_state['audio_buffer'].export(chunk_filename, format="wav")
                    
#                     try:
#                         async with aiohttp.ClientSession() as session:
#                             # Speech to text
#                             transcribed_text = await topic.speech_to_text(chunk_filename, username)
#                             output = filter_to_english(transcribed_text)
                            
#                             if not output:
#                                 logging.warning(f"Empty transcription for chunk {chunk_index}")
#                                 continue
                                
#                             # Send real-time transcription
#                             await websocket.send_json({"type": "transcription", "data": output})
#                             logging.info(f"[WS] {username} transcribed text: {output}")
                            
#                             # Grammar checking
#                             grammar_score = await topic.grammar_checking(transcribed_text)
#                             if grammar_score:
#                                 total_grammar_score += float(grammar_score)
#                                 grammar_score_count += 1

#                             # Emotion detection
#                             with open(chunk_filename, "rb") as f:
#                                 form = aiohttp.FormData()
#                                 form.add_field("file", f, filename=os.path.basename(chunk_filename), 
#                                              content_type="audio/wav")
#                                 async with session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as res:
#                                     emotion_data = await res.json()
#                                     emotion = emotion_data.get("emotion", "unknown")

#                             # Fluency scoring
#                             async with session.post(
#                                 f"{CPU_API_BASE}/fluency-score",
#                                 json={"text": output}
#                             ) as res:   
#                                 fluency_data = await res.json()
#                                 fluency = fluency_data.get("fluency", 0)

#                             # Pronunciation scoring
#                             with open(chunk_filename, "rb") as f:
#                                 form = aiohttp.FormData()
#                                 form.add_field("file", f, filename=os.path.basename(chunk_filename), 
#                                              content_type="audio/wav")
#                                 async with session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as res:
#                                     pron_data = await res.json()
#                                     pronunciation = pron_data.get("pronunciation", 0)

#                         # Update real-time stats
#                         topic.update_realtime_stats(fluency, pronunciation, emotion)
#                         text_output.append(output)

#                         chunk_result = {
#                             "chunk_index": chunk_index,
#                             "text": output,
#                             "emotion": emotion,
#                             "fluency": fluency,
#                             "pronunciation": pronunciation,
#                             "grammar_score": grammar_score,
#                             "silence_duration": session_state['silence_duration'],
#                             "file_path": chunk_filename
#                         }

#                         logging.info(f"[Chunk {chunk_index}] {chunk_result}")
#                         chunk_results.append(chunk_result)
                        
#                         # Reset session state
#                         session_state['audio_buffer'] = AudioSegment.empty()
#                         session_state['silence_duration'] = 0.0
#                         session_state['is_speaking'] = False
#                         chunk_index += 1
                        
#                     except Exception as e:
#                         logging.error(f"Error processing chunk {chunk_index}: {e}")
#                         # Clean up failed chunk file
#                         if os.path.exists(chunk_filename):
#                             os.remove(chunk_filename)
                
#     except WebSocketDisconnect:
#         logging.warning(f"[WS] {username} forcibly disconnected.")
#     except Exception as e:
#         logging.error(f"[WS] Unexpected error for {username}: {e}")
#     finally:
#         # Process any remaining audio in buffer
#         if len(session_state['audio_buffer']) > min_audio_length * 1000:
#             chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
#             session_state['audio_buffer'].export(chunk_filename, format="wav")
            
#             try:
#                 async with aiohttp.ClientSession() as session:
#                     transcribed_text = await topic.speech_to_text(chunk_filename, username)
#                     output = filter_to_english(transcribed_text)
                    
#                     if output:
#                         await websocket.send_json({"type": "transcription", "data": output})
#                         text_output.append(output)
                        
#                         grammar_score = await topic.grammar_checking(transcribed_text)
#                         if grammar_score:
#                             total_grammar_score += float(grammar_score)
#                             grammar_score_count += 1
                            
#                         chunk_result = {
#                             "chunk_index": chunk_index,
#                             "text": output,
#                             "emotion": "unknown",
#                             "fluency": 0,
#                             "pronunciation": 0,
#                             "grammar_score": grammar_score,
#                             "silence_duration": 0,
#                             "file_path": chunk_filename
#                         }
#                         chunk_results.append(chunk_result)
#             except Exception as e:
#                 logging.error(f"Error processing final chunk: {e}")
        
#         # Merge audio chunks and save transcript
#         loop = asyncio.get_event_loop()
#         await loop.run_in_executor(None, merge_chunks, chunk_results, final_output)

#         with open(transcript_path, "w", encoding="utf-8") as f:
#             f.write(" ".join(text_output).strip())

#         try:
#             # Get today's essays for this user
#             essays_ref = db.collection("essays").where("username", "==", username)
#             essays = essays_ref.stream()
#             today = datetime.now().date()
            
#             essay_list = []
#             for essay in essays:
#                 # Get the document data
#                 essay_data = essay.to_dict()
                
#                 # Get the creation timestamp from the document
#                 create_time = essay_data.get("timestamp")
#                 if not create_time:
#                     # If no timestamp field, use the document creation time
#                     create_time = essay.create_time
                
#                 # Convert to date for comparison
#                 if hasattr(create_time, 'date'):
#                     create_date = create_time.date()
#                 else:
#                     # Handle Firestore Timestamp objects
#                     create_date = create_time.to_python().date()
                
#                 if create_date == today:
#                     essay_list.append((essay.id, essay_data, create_time))
            
#             if essay_list:
#                 # Get the latest essay by creation time
#                 latest_essay_id, latest_essay_data, latest_create_time = max(
#                     essay_list, key=lambda x: x[2]
#                 )
                
#                 # Get a reference to the latest essay document
#                 essay_ref = db.collection("essays").document(latest_essay_id)

#                 average_scores = topic.get_average_realtime_scores()
#                 logging.info(f"Average Scores: {average_scores}")

#                 text_output_str = " ".join(text_output).strip()

#                 # Get the original text from the essay data
#                 original_text = latest_essay_data.get("content", "")
                
#                 # Calculate average grammar score
#                 if grammar_score_count > 0:
#                     total_grammar_score = round(total_grammar_score / grammar_score_count, 2)
#                 else:
#                     total_grammar_score = 0

#                 feedback_result = await final_listening__overall_scoring(
#                     average_scores, 
#                     text_output_str, 
#                     original_text,
#                     websocket,
#                     total_grammar_score
#                 )

#                 if "```json" in feedback_result:
#                     feedback_result = feedback_result.replace("```json", "")
#                 if "```" in feedback_result:
#                     feedback_result = feedback_result.replace("```", "")

#                 logging.info(f"Final Feedback: {feedback_result}")
                
#                 essay_ref.update({
#                     "chunks": chunk_results,
#                     "average_scores": average_scores,
#                     "overall_scoring": feedback_result
#                 })
#                 logging.info(f"Updated essay {latest_essay_id}")

#         except Exception as e:
#             logging.error(f"[Firestore Update Error] {e}")

#         finally:
#             # Clean up chunk files
#             for file in os.listdir(user_dir):
#                 if file.startswith("chunk_") and file.endswith(".wav"):
#                     try:
#                         os.remove(os.path.join(user_dir, file))
#                     except Exception as e:
#                         logging.warning(f"Failed to remove {file}: {e}")


# def merge_chunks(chunk_files, final_output):
#     logging.info("[Merge] Merging audio chunks...")
#     combined = AudioSegment.empty()

#     for chunk in chunk_files:
#         file_path = chunk.get("file_path")
#         if file_path and os.path.exists(file_path):
#             try:
#                 audio = AudioSegment.from_file(file_path, format="wav")
#                 combined += audio
#             except Exception as e:
#                 logging.warning(f"[Merge] Error reading {file_path}: {e}")
#         else:
#             logging.warning(f"[Merge] Skipping missing or invalid file: {file_path}")

#     combined.export(final_output, format="wav")
#     logging.info("[Merge] Merged audio file saved.")





@app.get("/get-tts-audio")
def get_tts_audio(username: str):
    folder = os.path.join("text_to_speech_audio_folder", username)
    file_path = os.path.join(folder, f"{username}_output.wav")

    timeout = 60
    poll_interval = 2
    waited = 0

    while waited < timeout:
        if os.path.exists(file_path):
            return FileResponse(file_path, media_type="audio/wav", filename=f"{username}_output.wav")
        time.sleep(poll_interval)
        waited += poll_interval
    raise HTTPException(status_code=408, detail="Audio file not generated within 1 minute.")



SUPPORTED_IMAGE_FORMATS = [".png", ".jpg", ".jpeg"]
SUPPORTED_TEXT_FORMATS = [".txt"]
SUPPORTED_PDF_FORMATS = [".pdf"]
SUPPORTED_DOC_FORMATS = [".docx"]
SUPPORTED_PPT_FORMATS = [".pptx"]
SUPPORTED_XLS_FORMATS = [".xlsx"]


genai.configure(api_key=GOOGLE_API_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-pro")
pc = Pinecone(api_key=PINECONE_API_KEY)

if PINECONE_INDEX_NAME not in pc.list_indexes().names():
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=1024,
        metric="cosine",
        pods=1,
        pod_type="p1.x1"
    )
    logging.info(f"Created new Pinecone index {PINECONE_INDEX_NAME} with dimension 1024")

index = pc.Index(PINECONE_INDEX_NAME)

embedding_model = HuggingFaceEmbeddings(
    model_name=EMBEDDING_MODEL_NAME,
    model_kwargs={'device': 'cuda' if torch.cuda.is_available() else 'cpu'}
)

executor = ThreadPoolExecutor(max_workers=8)

async def run_in_threadpool(func, *args, **kwargs):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(executor, lambda: func(*args, **kwargs))


def render_text_to_image(text: str, width=800, font_size=18) -> Image.Image:
    font = ImageFont.load_default()
    lines = []
    dummy_img = Image.new("RGB", (width, 1000))
    draw = ImageDraw.Draw(dummy_img)

    words = text.split()
    line = ""
    for word in words:
        test_line = f"{line} {word}".strip()
        bbox = draw.textbbox((0, 0), test_line, font=font)
        w = bbox[2] - bbox[0]
        if w < width - 40:
            line = test_line
        else:
            lines.append(line)
            line = word
    lines.append(line)

    height = font_size * len(lines) + 50
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)
    y = 20
    for line in lines:
        draw.text((20, y), line, font=font, fill="black")
        y += font_size
    return img

def file_to_images(file_path: str) -> List[Image.Image]:
    ext = os.path.splitext(file_path)[1].lower()
    images = []

    if ext in SUPPORTED_PDF_FORMATS:
        with ThreadPoolExecutor() as pdf_executor:
            futures = []
            chunk_size = 10
            images = convert_from_path(file_path, dpi=200, thread_count=4)

    elif ext in SUPPORTED_IMAGE_FORMATS:
        images = [Image.open(file_path)]

    elif ext in SUPPORTED_DOC_FORMATS:
        try:
            text = docx2txt.process(file_path)
            img = render_text_to_image(text)
            images = [img]
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"DOCX to image failed: {e}")

    elif ext in SUPPORTED_PPT_FORMATS:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            img = render_text_to_image(text)
            images.append(img)

    elif ext in SUPPORTED_XLS_FORMATS:
        try:
            excel = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in excel.items():
                text = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                img = render_text_to_image(text)
                images.append(img)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"XLSX to image failed: {e}")

    elif ext in SUPPORTED_TEXT_FORMATS:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        img = render_text_to_image(content)
        images = [img]

    else:
        raise HTTPException(status_code=400, detail="Unsupported file format")

    return images


async def process_single_image(filename,username,image: Image.Image, idx: int,image_path) -> str:
    try:
        today_date = datetime.now().strftime("%Y-%m-%d")
        base_dir = f"diagrams/{username}_{today_date}?extracted_images"
        
        response = await run_in_threadpool(
            gemini_model.generate_content,
            [
                image,
                "1. Extract all visible text from this image.\n"
                "2. If the image contains diagrams, illustrations, or charts, summarize what they represent in 2-3 lines.\n"
                "Return the result with clear separation: first the text, then the image summary (if any)."
            ]
        )
        text = response.text.strip()

        result = subprocess.run(
            [python311_path, script_path, username, image_path, str(idx), filename],
            capture_output=True,
            text=True
        )
        try:
            output = json.loads(result.stdout)
            print(f"Page {idx + 1} Output:", output["result"])
        except json.JSONDecodeError:
            print(f"Error decoding output for page {idx + 1}:", result.stdout)

        if result.stderr:
            print(f"Error from script on page {idx + 1}:", result.stderr)

        return f"\n\n--- Page/Image {idx + 1} ---\n{text}"
    except Exception as e:
        logging.error(f"OCR failed on image {idx + 1}: {e}")
        return f"\n\n--- Page/Image {idx + 1} FAILED ---"

async def extract_text_parallel(filename,username,file_path: str, timeout_per_page: int = 240) -> str:
    images = await run_in_threadpool(file_to_images, file_path)
    if not images:
        return ""
    
    all_text = ""
    tasks = []

    filename = filename.replace(" ","_")
    filename = filename.split(".")[0]
    output_image_dir = f"/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi/temp_images/{username}/{filename}"
    if os.path.exists(output_image_dir):
        if os.listdir(output_image_dir):
            shutil.rmtree(output_image_dir)

    os.makedirs(output_image_dir, exist_ok=True)

    
    
    for idx, image in enumerate(images):
        image_path = os.path.join(output_image_dir, f"page_{idx + 1}.jpg")
        image.save(image_path, "JPEG")
        task = asyncio.create_task(
            asyncio.wait_for(
                process_single_image(filename,username,image, idx,image_path),
                timeout=timeout_per_page
            )
        )
        tasks.append(task)
    
    for i, task in enumerate(asyncio.as_completed(tasks)):
        try:
            page_text = await task
            all_text += page_text
        except asyncio.TimeoutError:
            logging.warning(f"Timeout processing page {i + 1}")
            all_text += f"\n\n--- Page/Image {i + 1} TIMEOUT ---"
        except Exception as e:
            logging.error(f"Error processing page {i + 1}: {e}")
            all_text += f"\n\n--- Page/Image {i + 1} ERROR ---"
    
    return all_text

async def extract_text_with_retry(filename,username,file_path: str, timeout=240, max_retries=3) -> tuple[str, bool]:
    last_error = None
    timeout_occurred = False
    
    for attempt in range(1, max_retries + 2):
        try:
            task_start = time.time()
            
            task = asyncio.create_task(extract_text_parallel(filename,username,file_path))
            
            try:
                result = await asyncio.wait_for(task, timeout=timeout)
                elapsed = time.time() - task_start
                logging.info(f"OCR succeeded in attempt {attempt} ({elapsed:.2f}s)")
                return result, timeout_occurred
                
            except asyncio.TimeoutError:
                task.cancel()
                elapsed = time.time() - task_start
                logging.warning(
                    f"OCR timeout in attempt {attempt} after {elapsed:.2f}s "
                    f"(Timeout setting: {timeout}s)"
                )
                last_error = f"Timeout after {timeout}s"
                timeout_occurred = True
                
                if attempt <= max_retries:
                    backoff = min(2 ** attempt, 10)
                    await asyncio.sleep(backoff)
                    
        except Exception as e:
            last_error = str(e)
            logging.error(f"OCR attempt {attempt} failed: {last_error}")
            if attempt <= max_retries:
                await asyncio.sleep(1)

    raise HTTPException(
        status_code=504,
        detail=f"OCR failed after {max_retries + 1} attempts. Last error: {last_error}"
    )

@app.post("/upload/")
async def upload_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    student_class: str = Form(...),
    subject: str = Form(...),
    curriculum: str = Form(...),
    username : str = Form(...),
    
):

    file_path = None
    try:
        start_time = time.time()
        folder = f"uploads/{curriculum}/{student_class}/{subject}"
        os.makedirs(folder, exist_ok=True)
        file_path = os.path.join(folder, file.filename)

        filename = file.filename

        with open(file_path, "wb") as buffer:
            while chunk := await file.read(1024 * 1024):
                buffer.write(chunk)
        logging.info(f"File saved in {time.time()-start_time:.2f}s")

        extract_start = time.time()
        try:
            extracted_text, timeout_occurred = await extract_text_with_retry(
                filename,
                username,
                file_path,
                timeout=240,
                max_retries=2,
            )
            logging.info(f"Text extracted in {time.time()-extract_start:.2f}s")
            
            if not extracted_text.strip():
                raise HTTPException(
                    status_code=422,
                    detail="No text could be extracted from the file"
                )
                
        except HTTPException as e:
            raise e
        except Exception as e:
            logging.error(f"OCR processing failed: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail="Failed to process document content"
            )

        namespace = f"{curriculum}_{student_class}_{subject}"
        
        existing_entries = index.query(
            vector=embedding_model.embed_query(extracted_text[:1000]),
            top_k=1,
            filter={
                "filename": {"$eq": file.filename},
                "type": {"$eq": "document"}
            },
            namespace=namespace
        )
        
        if existing_entries.matches:
            return {
                "status": "duplicate",
                "message": "File already exists in vector database",
                "existing_id": existing_entries.matches[0].id
            }

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(extracted_text)
        
        metadatas = [{
            "curriculum": curriculum,
            "student_class": student_class,
            "subject": subject,
            "filename": file.filename,
            "type": "document",
            "chunk_idx": i,
            "text": chunk[:500]
        } for i, chunk in enumerate(chunks)]

        background_tasks.add_task(
            store_in_vector_db,
            chunks=chunks,
            metadatas=metadatas,
            namespace=namespace
        )

        return {
            "status": "success",
            "filename": file.filename,
            "message": "File processed and queued for vector storage",
            "processing_times": {
                "file_save": f"{time.time()-start_time:.2f}s",
                "text_extraction": f"{time.time()-extract_start:.2f}s"
            },
            "timeout_occurred": timeout_occurred,
            "chunk_count": len(chunks),
            "namespace": namespace
        }

    except HTTPException as e:
        raise e
    except Exception as e:
        logging.error(f"Upload failed: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail="Internal server error during file processing"
        )
    finally:
        if file_path and os.path.exists(file_path):
            try:
                os.remove(file_path)
                logging.info("Temporary file cleaned up")
            except Exception as e:
                logging.warning(f"Failed to delete temp file: {str(e)}")

def store_in_vector_db(chunks: List[str], metadatas: List[Dict], namespace: str):
    try:
        logging.info(f"Starting vector storage for {len(chunks)} chunks in namespace {namespace}")
        
        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text",
            namespace=namespace
        )
        
        batch_size = 50 
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i:i + batch_size]
            batch_metadatas = metadatas[i:i + batch_size]
            
            try:
                embeddings = embedding_model.embed_documents(batch_chunks)
                
                records = []
                for j, (chunk, metadata) in enumerate(zip(batch_chunks, batch_metadatas)):
                    record = {
                        "id": f"{namespace}_{i+j}_{uuid.uuid4().hex[:8]}",
                        "values": embeddings[j],
                        "metadata": metadata
                    }
                    records.append(record)
                
                vectorstore._index.upsert(vectors=records, namespace=namespace)
                logging.info(f"Upserted batch {i//batch_size + 1} with {len(records)} vectors")
                
            except Exception as e:
                logging.error(f"Failed to process batch {i//batch_size + 1}: {str(e)}")
                continue
        
        logging.info(f"Completed vector storage for namespace {namespace}")
        
    except Exception as e:
        logging.error(f"Vector storage failed: {str(e)}", exc_info=True)


BASE_DIR = Path("/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi")

class ChatResponse(TypedDict):
    question: str
    answer: Annotated[str, "Do not include URLs, only the valuable answer."]
    diagram_urls: Annotated[str, "Provide exactly one relevant diagram URL."]
    source_documents: str

model_name = OllamaLLM(model="mistral")

def parse_llm_response(response_text: str):
    """Parse the LLM response which should be in JSON format"""
    try:
        start_idx = response_text.find('{')
        end_idx = response_text.rfind('}') + 1
        json_str = response_text[start_idx:end_idx]
        return json.loads(json_str)
    except (json.JSONDecodeError, ValueError) as e:
        logging.error(f"Failed to parse LLM response: {e}")
        return {
            "question": "",
            "answer": response_text,
            "url": ""
        }   

HISTORY_DIR = "/chat_histroy_tmp"
MAX_HISTORY = 10

HISTORY_DIR = os.path.join(os.getcwd(), "chat_histroy_tmp")
os.makedirs(HISTORY_DIR, exist_ok=True)

def load_chat_history(file_path):
    """Load existing chat history or return empty list."""
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except json.JSONDecodeError:
            logging.warning(f"Corrupted history file: {file_path}, resetting.")
            return []
    return []

def save_chat_history(file_path, history):
    """Save chat history safely (atomic write)."""
    tmp_path = f"{file_path}.tmp"
    with open(tmp_path, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=2)
    os.replace(tmp_path, file_path)

def cleanup_old_history(username, student_class, subject, curriculum, today_date):
    """Delete all history files for the user that are not from today."""
    pattern = f"{username}_{student_class}_{subject}_{curriculum}_*.json"
    old_files = glob.glob(os.path.join(HISTORY_DIR, pattern))
    for file_path in old_files:
        if today_date not in os.path.basename(file_path):
            try:
                os.remove(file_path)
                logging.info(f"Deleted old history file: {file_path}")
            except Exception as e:
                logging.error(f"Failed to delete old history file {file_path}: {e}")

@app.post("/chat/", response_model=ChatResponse)
async def chat(request: ChatRequest):
    try:
        today_date = datetime.now().strftime("%Y-%m-%d")
        chat_student_class = request.student_class.strip().replace(" ", "_")
        cleanup_old_history(
            request.username.strip(),
            chat_student_class,
            request.subject.strip(),
            request.curriculum.strip(),
            today_date
        )

        history_filename = (
            f"{request.username.strip()}_"
            f"{chat_student_class}_"
            f"{request.subject.strip()}_"
            f"{request.curriculum.strip()}_"
            f"{today_date}.json"
        )
        history_file_path = os.path.join(HISTORY_DIR, history_filename)

        chat_history = load_chat_history(history_file_path)

        namespace = f"{request.curriculum.strip()}_{request.student_class.strip()}_{request.subject.strip()}"
        base_dir = f"/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi/diagrams/{request.username.strip()}_{today_date}/extracted_images"
        image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.gif', '*.bmp', '*.tiff']
        image_files = []
        for ext in image_extensions:
            image_files.extend(glob.glob(os.path.join(base_dir, ext)))

        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text",
            namespace=namespace
        )
        retriever = vectorstore.as_retriever(search_kwargs={
            "k": 8,
            "filter": {
                "subject": request.subject.strip(),
                "curriculum": request.curriculum.strip(),
                "student_class": request.student_class.strip(),
            }
        })

        retrieved_docs = retriever.invoke(request.question.strip())
        retrieved_docs_content = " ".join(doc.page_content for doc in retrieved_docs)

        history_context = "\n".join(
            f"Q: {h['question']}\nA: {h['answer']}" for h in chat_history[-MAX_HISTORY:]
        )

        prompt_template = PromptTemplate(
            template="""
            ROLE: You are an expert educator providing clear, concise answers to students.

            INSTRUCTIONS:
            1. Use ONLY the provided context and chat history to answer the question
            2. Analyze the query type:
            - If it's a normal conversational query, use history context and reply directly
            - If it's a specific question, use the provided context to answer
            3. Image selection criteria:
            - Use exactly ONE relevant image ONLY if relevance exceeds 50%
            - Image path must be EXACTLY one of the provided paths
            - If no suitable image exists, use empty string for url

            INPUT DATA:
            Chat History: {history}

            Context: {context}

            Available Images: {image_files}

            Question: {question}

            OUTPUT REQUIREMENTS:
            - Respond in STRICT JSON format ONLY
            - The answer must NOT contain any URLs or external references
            - Follow this EXACT structure:
            {{
                "question": "the original asked question",
                "answer": "your comprehensive answer without any URLs",
                "url": "exact matching image path or empty string"
            }}

            VALIDATION RULES:
            1. Before responding, verify:
            - Have I checked both history and context appropriately?
            - Is the image selected truly relevant (>65% match)?
            - Does my answer contain no URLs or external references?
            - Is the output in exact JSON format required?

            2. Error prevention:
            - If question is unclear, ask for clarification in the answer field
            - If no relevant context exists, state this clearly in the answer
            - Never invent information beyond what's provided

            RESPONSE:
            """,
                input_variables=["history", "question", "image_files", "context"]
        )
        qa_chain = LLMChain(llm=model_name, prompt=prompt_template)

        result = qa_chain.invoke({
            "history": history_context,
            "question": request.question.strip(),
            "image_files": "\n".join(image_files) if image_files else "No images available",
            "context": retrieved_docs_content
        })

        parsed_response = parse_llm_response(result["text"])
        response_url = parsed_response.get("url", "")

        if response_url == "empty string" or response_url == "empty_string" or response_url == "":
            response_url = ""

        elif response_url.lstrip().startswith("/diagrams"):
            response_url = response_url

        elif response_url.lstrip().startswith("/home"):
            response_url = response_url.split("/diagrams")[1]
            response_url = f"/diagrams{response_url}"
            
        else:
            response_url = f"/diagrams/{request.username.strip()}_{today_date}/extracted_images/{response_url}.png"

        if response_url.endswith(".png.png"):
             response_url = response_url.removesuffix(".png.png") + ".png"

        full_answer = parsed_response.get("answer","")
        
        full_answer = re.sub(r"(/home\S+\.(?:png|jpg|jpeg|gif|bmp|tiff))", "", full_answer, flags=re.IGNORECASE)
        full_answer = re.sub(r"(/diagrams\S+\.(?:png|jpg|jpeg|gif|bmp|tiff))", "", full_answer, flags=re.IGNORECASE)
        full_answer = re.sub(r"\s+", " ", full_answer).strip()

        chat_history.append({"question": request.question.strip(), "answer": full_answer})
        if len(chat_history) > MAX_HISTORY:
            chat_history = chat_history[-MAX_HISTORY:]
        save_chat_history(history_file_path, chat_history)

        return ChatResponse(
            question=parsed_response.get("question", request.question.strip()),
            answer=full_answer,
            diagram_urls=response_url,
            source_documents=retrieved_docs_content
        )

    except Exception as e:
        logging.error(f"Chat error: {str(e)}", exc_info=True)
        return ChatResponse(
            question=request.question,
            answer="An error occurred while processing your request.",
            diagram_urls="",
            source_documents="",
        )


@app.get("/view-image/{full_path:path}")
def view_image(full_path: str):
    image_path = (BASE_DIR / full_path.lstrip("/")).resolve()

    if not str(image_path).startswith(str(BASE_DIR)):
        raise HTTPException(status_code=400, detail="Invalid file path")

    if not image_path.exists():
        raise HTTPException(status_code=404, detail="Image not found")

    return FileResponse(image_path)


@app.get("/health")
def welcome_page():
    return {"Message": "Welcome the ai speech module page."}

chat_history = []

model = ChatOllama(
            model="mistral", 
            model_kwargs={"temperature": 0.8}
        )

async def system_message(topic, mood, student_class, level) -> SystemMessage:
    parser = StrOutputParser()

    prompt_template = PromptTemplate(template="""
    You are a friendly, knowledgeable teaching assistant. Your purpose is to:
    1. Introduce the topic in a friendly manner
    2. Answer questions conversationally
    3. Never reveal internal project details
    4. Keep responses under 100 words
                                     
    Topic: {topic}
    Mood: {mood}
    Student Class: {student_class}
    Level: {level}
    5. Introduce on understanding the topic always.""",
    input_variables=["topic", "mood", "student_class", "level"])

    chain = prompt_template | model | parser
    result = await chain.ainvoke({
        "topic": topic,
        "mood": mood,
        "student_class": student_class,
        "level": level
    })
    return result

async def initialize_essay_document(username, student_topic, student_class, mood, accent, chat_history):
    """Initialize a new essay document in Firestore"""
    try:
        serializable_chat_history = []
        for message in chat_history:
            if isinstance(message, (AIMessage, SystemMessage, HumanMessage)):
                message_dict = {
                    'type': 'chat_message',
                    'content': message.content,
                    'message_type': 'system' if isinstance(message, SystemMessage) else 
                                   'ai' if isinstance(message, AIMessage) else 
                                   'human'
                }
                serializable_chat_history.append(message_dict)
            else:
                serializable_chat_history.append(str(message))

        essay_data = {
            "username": username,
            "topic": student_topic,
            "student_class": student_class,
            "mood": mood,
            "accent": accent,
            "chat_history": serializable_chat_history,
            "created_at": datetime.now(),
            "status": "in_progress",
            "chunks": [],
            "average_scores": {}
        }

        _, doc_ref = db.collection("essays").add(essay_data)
        return doc_ref.id
        
    except Exception as e:
        logging.error(f"Failed to initialize essay document: {str(e)}", exc_info=True)
        raise

from collections import deque

feedback_queue = deque()
feedback_event = asyncio.Event()


async def feedback_processor():
    while True:
        await feedback_event.wait()
        while feedback_queue:
            websocket, feedback = feedback_queue.popleft()
            try:
                await websocket.send_json({
                    "type": "feedback",
                    "data": feedback
                })
            except Exception as e:
                logging.error(f"Failed to send feedback: {str(e)}")
        feedback_event.clear()

@app.on_event("startup")
async def startup_event():
    asyncio.create_task(feedback_processor())


async def finalize_listening_session(result, transcribed_text, websocket):
    try:
        fluency = result.get("fluency")
        pronunciation = result.get("pronunciation")
        emotion = result.get("emotion")
        grammar_score = result.get("grammar")

        prompt_template = PromptTemplate(template = """
            You are an expert English teacher providing comprehensive feedback on a student's listening comprehension and speaking response.

            CONTEXT ANALYSIS:
            The student has completed a listening exercise where they heard content and are now providing their spoken response.
            Analyze both what they said AND how they said it.

            STUDENT'S COMPLETE RESPONSE:
            "{transcribed_text}"

            PERFORMANCE METRICS:
            - Average Fluency: {fluency}/10 (smoothness and flow of speech)
            - Average Pronunciation: {pronunciation}/10 (clarity and accuracy of word pronunciation)
            - Dominant Emotion: {emotion} (emotional tone during delivery)
            - Grammar Score: {grammar_score}/10 (grammatical accuracy)

            ANALYSIS GUIDELINES:
            1. For content analysis: Consider what the response reveals about listening comprehension
            2. For speaking performance: Focus on delivery quality regardless of content accuracy
            3. Be specific: Provide exact examples from the student's response
            4. Be constructive: Frame feedback positively with actionable suggestions
            5. Consider context: The student heard content before speaking, so assess recall and understanding

            DETAILED ANALYSIS REQUIRED:
            Strict Provide comprehensive feedback in JSON format with these specific sections:

            {{
                "feedback": {{
                    "content_understanding": {{
                        "score": "X/10",
                        "explanation": "How well did the student demonstrate understanding of the listened content? Analyze coherence, relevance, and depth of response.",
                        "evidence": "Specific phrases or ideas that show comprehension level"
                    }},
                    "detail_retention": {{
                        "score": "X/10", 
                        "explanation": "How well did they retain and articulate specific details from what they heard?",
                        "specifics_mentioned": ["Key details recalled accurately"],
                        "approximations": ["Details recalled with minor inaccuracies"]
                    }},
                    "key_points_covered": ["Main themes or ideas the student successfully captured"],
                    "potential_missed_opportunities": ["Areas where more depth or detail could have been added based on typical listening content"]
                }},
                "speaking_performance": {{
                    "fluency_assessment": {{
                        "score": {fluency},
                        "analysis": "Detailed analysis of speech flow, hesitations, and natural rhythm with specific timestamp examples",
                        "strengths": "Smooth segments and effective pacing examples",
                        "improvement_areas": "Specific moments where fluency broke down"
                    }},
                    "pronunciation_assessment": {{
                        "score": {pronunciation},
                        "analysis": "Comprehensive pronunciation evaluation with word-specific feedback",
                        "well_prnounced_words": ["List of correctly pronounced challenging words"],
                        "needs_work_words": ["Problematic words with phonetic guidance: word→[correct pronunciation]"]
                    }},
                    "grammar_assessment": {{
                        "score": {grammar_score},
                        "analysis": "Detailed grammatical analysis with error categorization",
                        "error_examples": ["Incorrect: 'she go' → Correct: 'she goes'"],
                        "complexity_level": "Assessment of sentence structures attempted"
                    }},
                    "vocabulary_usage": {{
                        "assessment": "Evaluation of word choice appropriateness and variety",
                        "effective_vocabulary": ["Sophisticated or context-appropriate words used well"],
                        "vocabulary_opportunities": ["Where more precise vocabulary could enhance meaning"]
                    }},
                    "speaking_clarity": {{
                        "overall_rating": "Assessment of how understandable the response was",
                        "clarity_factors": ["Pacing", "Articulation", "Volume consistency", "Phrasing"]
                    }}
                }},
                "technical_metrics": {{
                    "speaking_rate_analysis": "Assessment of words per minute and pace variation appropriateness",
                    "pause_analysis": "Analysis of pause placement: effective pauses vs. hesitation pauses",
                    "filler_word_usage": "Count and analysis of filler words (um, ah, like) with impact assessment",
                    "prosody_evaluation": "Intonation patterns, stress placement, and rhythmic flow quality"
                }},
                "detailed_suggestions": [
                    {{
                        "suggestion": "Specific actionable suggestion 1",
                        "example": "Concrete example from student's speech",
                        "improved_version": "How to implement the suggestion"
                    }},
                    {{
                        "suggestion": "Specific actionable suggestion 2", 
                        "example": "Concrete example from student's speech",
                        "improved_version": "How to implement the suggestion"
                    }},
                    {{
                        "suggestion": "Specific actionable suggestion 3",
                        "example": "Concrete example from student's speech",
                        "improved_version": "How to implement the suggestion"
                    }},
                    {{
                        "suggestion": "Specific actionable suggestion 4",
                        "example": "Concrete example from student's speech",
                        "improved_version": "How to implement the suggestion"
                    }},
                    {{
                        "suggestion": "Specific actionable suggestion 5",
                        "example": "Concrete example from student's speech",
                        "improved_version": "How to implement the suggestion"
                    }}
                ],
                "strengths": [
                    {{
                        "strength": "Specific strength 1",
                        "evidence": "Exact phrase or moment demonstrating this strength without inside ' ' ",
                        "impact": "Why this strength is effective"
                    }},
                    {{
                        "strength": "Specific strength 2",
                        "evidence": "Exact phrase or moment demonstrating this strength without inside ' ' ",
                        "impact": "Why this strength is effective"
                    }},
                    {{
                        "strength": "Specific strength 3",
                        "evidence": "Exact phrase or moment demonstrating this strength without inside ' ' ",
                        "impact": "Why this strength is effective"
                    }}
                ],
                "practice_recommendations": [
                    {{
                        "activity": "Specific practice activity 1",
                        "purpose": "What skill this develops",
                        "frequency": "Recommended practice schedule"
                    }},
                    {{
                        "activity": "Specific practice activity 2",
                        "purpose": "What skill this develops",
                        "frequency": "Recommended practice schedule"
                    }},
                    {{
                        "activity": "Specific practice activity 3", 
                        "purpose": "What skill this develops",
                        "frequency": "Recommended practice schedule"
                    }}
                ],
                "overall_scores": {{
                    "fluency": {fluency},
                    "pronunciation": {pronunciation},
                    "grammar": {grammar_score},
                    "emotion": "{emotion}",
                    "comprehension_quality": "X/10 (based on content analysis)"
                }},
                "improvement_priority": {{
                    "area": "The single most important area to focus on next",
                    "reason": "Why this area is most critical for improvement",
                    "immediate_action": "First step to address this priority"
                }},
                "encouragement": {{
                    "progress_highlight": "Specific improvement or strength noticed",
                    "motivational_message": "Personalized encouraging message",
                    "growth_potential": "What they can achieve with continued practice"
                }}
            }}

            CRITICAL: Ensure all feedback is specific, actionable, and includes concrete examples from the student's actual response.
            """,
            input_variables=["transcribed_text", "fluency", "pronunciation", "emotion", "grammar_score"]
            )
        parser = StrOutputParser()
        chain = prompt_template | model | parser
        feedback = await chain.ainvoke({
            "transcribed_text": transcribed_text,
            "fluency": fluency,
            "pronunciation": pronunciation,
            "emotion": emotion,
            "grammar_score": grammar_score,
            "fluency": fluency,
            "pronunciation": pronunciation,
            "emotion": emotion,
            "grammar_score": grammar_score,
        })

        feedback_queue.append((websocket, feedback))
        feedback_event.set()
        
        return feedback

    except Exception as e:
        logging.info(f"Session finalization error: {e}")


import contextlib

async def _safe_finalize_feedback(
    config: dict,
    temp_dir: str,
    text: str,
    websocket: WebSocket,
    session_state: dict,
    username: str,
    buffer_filename: str,
    topic: Topic,
    http_session: aiohttp.ClientSession,
    final_feedback_result: list
):
    try:
        audio_buf = session_state.get('audio_buffer', AudioSegment.empty())
        grammar_score = await topic.grammar_checking(text)
        
        final_output = os.path.join(temp_dir, f"{username}_output.wav")
        if os.path.exists(final_output):
            existing_audio = AudioSegment.from_wav(final_output)
            combined_audio = existing_audio + audio_buf
        else:
            combined_audio = audio_buf
        combined_audio.export(final_output, format="wav")

        transcript_path = os.path.join(temp_dir, f"{username}_transcript.txt")
        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(" ".join(session_state.get('text_output', [])).strip())

        clean_text = text.lower().strip()
        result = {}
        if (len(text.split()) >= config['min_utterance_length']):

            with open(buffer_filename, "rb") as f:
                form = aiohttp.FormData()
                form.add_field("file", f, filename=os.path.basename(buffer_filename), content_type="audio/wav")
                async with http_session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as res:
                    emotion_data = await res.json()
                    result["emotion"] = (emotion_data or {}).get("emotion")

            async with http_session.post(f"{CPU_API_BASE}/fluency-score", json={"text": text}) as res:
                fluency_data = await res.json()
                result["fluency"] = (fluency_data or {}).get("fluency")

            with open(buffer_filename, "rb") as f:
                form = aiohttp.FormData()
                form.add_field("file", f, filename=os.path.basename(buffer_filename), content_type="audio/wav")
                async with http_session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as res:
                    pron_data = await res.json()
                    result["pronunciation"] = (pron_data or {}).get("pronunciation")

        result["grammar"] = grammar_score

        config["fluency_score"] += float(result.get("fluency", 0))
        config["pronunciation_score"] += float(result.get("pronunciation", 0))
        config["emotion_score"] = result.get("emotion", 0)
        config["grammar_score"] += float(result.get("grammar", 0))
        config["count"] += 1

        topic.update_realtime_stats(result.get("fluency"), result.get("pronunciation"), result.get("emotion"))

        session_state.setdefault('chunk_results', []).append({
            "chunk_index": session_state.get('chunk_index'),
            "text": text,
            "emotion": result.get("emotion"),
            "fluency": result.get("fluency"),
            "pronunciation": result.get("pronunciation"),
            "grammar": grammar_score
        })

        if None in result.values():
            logging.warning(f"Incomplete analysis results: {result}")

        logging.info(f"result is --------------->{result}")

        feedback_result = await finalize_listening_session(result, text, websocket)
        final_feedback_result.append(feedback_result)
        session_state["last_feedback"] = feedback_result
        logging.info(f"this is feedback of chunk : {feedback_result}")

    except Exception:
        logging.exception("finalize_listening_session failed")


async def overall_scoring_by_speech_module_test(
    final_feedback_result: list, 
    session_state: dict, 
    topic: "Topic",
    config: dict,
    essay_id: str,
    store_in_db: bool = True
) -> Optional[Dict[str, Any]]:
    logging.info("coming inside the overall scoring funtion -------------------------->>>>>>>>>>>>")
    try:
        count = config.get("count", 0)
        if count == 0:
            logging.warning("No audio chunks were processed, skipping overall scoring")
            return None

        fluency_score = config.get("fluency_score", 0) / count
        pronunciation_score = config.get("pronunciation_score", 0) / count
        grammar_score = config.get("grammar_score", 0) / count
        emotion_score = config.get("emotion_score", "neutral")
        
        logging.info(f"Overall scores - Fluency: {fluency_score:.2f}, Pronunciation: {pronunciation_score:.2f}, Grammar: {grammar_score:.2f}, Emotion: {emotion_score}")

        context_text = "No feedback data available"
        if final_feedback_result:
            try:
                context_text = json.dumps(final_feedback_result, ensure_ascii=False, indent=2)
            except Exception as e:
                logging.warning(f"Failed to serialize feedback results: {e}")
                context_text = str(final_feedback_result)

        prompt_template = PromptTemplate(
            template="""
            You are a top English teacher providing comprehensive feedback on a student's speaking performance.

            CONTEXT ANALYSIS:
            You are analyzing aggregated feedback from multiple speech chunks from a listening comprehension exercise.
            The student heard content and provided spoken responses that were chunked and analyzed separately.

            AGGREGATED FEEDBACK DATA FROM ALL CHUNKS:
            {final_feedback_result}

            PERFORMANCE SUMMARY:
            - Average Fluency: {fluency}/10 (across all speech segments)
            - Average Pronunciation: {pronunciation}/10 (across all words spoken)
            - Average Grammar: {grammar_score}/10 (overall grammatical accuracy)
            - Dominant Emotion: {emotion} (prevailing emotional tone)

            ANALYSIS GUIDELINES:
            1. SYNTHESIZE PATTERNS: Identify consistent strengths and weaknesses across all chunks
            2. CONTEXT AWARENESS: Remember this was a listening exercise - assess comprehension, not just speaking
            3. SPECIFIC EXAMPLES: Reference actual content from the feedback chunks
            4. PROPORTIONAL FEEDBACK: Weight feedback based on frequency of occurrences
            5. ACTIONABLE INSIGHTS: Provide concrete, implementable suggestions

            CRITICAL CONSIDERATIONS:
            - If chunks show inconsistent performance, analyze why (fatigue, complexity, etc.)
            - For listening comprehension: assess how well student captured main ideas vs details
            - Consider speaking performance independent of content accuracy
            - Account for natural speech patterns vs systematic errors

            REQUIRED OUTPUT FORMAT (STRICT JSON):
            {{
                "feedback": {{
                    "content_understanding": {{
                        "score": "X/10",
                        "explanation": "How well did the student understand the main ideas across all responses? Analyze consistency and depth.",
                        "pattern_analysis": "Were there patterns in comprehension across chunks?",
                        "evidence": "Specific examples from aggregated feedback showing comprehension level"
                    }},
                    "detail_retention": {{
                        "score": "X/10", 
                        "explanation": "How consistently did they retain specific details across different segments?",
                        "key_points_covered": ["Main points student consistently mentioned across chunks"],
                        "missed_points": ["Important points consistently missed across chunks"],
                        "improvement_trend": "Did detail retention improve/worsen over time?",
                    }},
                    "comprehension_consistency": "Analysis of how understanding varied across different parts of the content",
                }},
                "speaking_performance": {{
                    "fluency_assessment": {{
                        "score": {fluency},
                        "analysis": "Detailed fluency analysis across all chunks. Identify patterns: consistent issues or situational factors.",
                        "consistency": "Was fluency stable or variable across the session?",
                        "specific_examples": ["Timestamps or contexts where fluency excelled/failed"],
                    }},
                    "pronunciation_assessment": {{
                        "score": {pronunciation},
                        "analysis": "Pronunciation patterns across all speech. Identify systematic vs occasional errors.",
                        "consistent_strengths": ["Words consistently pronounced correctly"],
                        "systematic_errors": ["Words consistently mispronounced with corrections"],
                        "situational_errors": ["Pronunciation issues in specific contexts"],
                    }},
                    "grammar_assessment": {{
                        "score": {grammar_score},
                        "analysis": "Grammar accuracy patterns. Distinguish between slips and systematic errors.",
                        "error_patterns": ["Recurring grammatical mistakes across chunks"],
                        "complexity_attempts": "Level of grammatical structures attempted consistently",
                        "self_correction": "Evidence of self-monitoring and correction",
                    }},
                    "vocabulary_usage": {{
                        "assessment": "Vocabulary range and appropriateness across entire session",
                        "consistency": "Was vocabulary usage consistent or variable?",
                        "appropriateness": "How well vocabulary matched the content level throughout"
                    }},
                    "speaking_clarity": {{
                        "overall_rating": "Consistency of clarity across different speaking segments",
                        "clarity_trend": "Did clarity improve or deteriorate over time?"
                    }}
                }},
                "technical_metrics": {{
                    "speaking_rate_consistency": "Analysis of pace variation across the session",
                    "pause_patterns": "Were pauses strategic or hesitant? Pattern analysis",
                    "filler_word_patterns": "Analysis of filler word usage frequency and contexts",
                    "prosody_consistency": "How consistent were intonation and stress patterns?"
                }},
                "detailed_suggestions": [
                    {{
                        "suggestion": "Specific suggestion 1 addressing most frequent issue",
                        "based_on": "Which feedback chunks support this suggestion",
                        "implementation": "Concrete steps to implement",
                        "expected_impact": "How this will improve performance"
                    }},
                    {{
                        "suggestion": "Specific suggestion 2 addressing pattern from chunks", 
                        "based_on": "Which feedback chunks support this suggestion",
                        "implementation": "Concrete steps to implement",
                        "expected_impact": "How this will improve performance"
                    }},
                    {{
                        "suggestion": "Specific suggestion 3 for consistent strength development",
                        "based_on": "Which feedback chunks support this suggestion",
                        "implementation": "Concrete steps to implement",
                        "expected_impact": "How this will improve performance"
                    }},
                    {{
                        "suggestion": "Specific suggestion 4 for situational improvement",
                        "based_on": "Which feedback chunks support this suggestion",
                        "implementation": "Concrete steps to implement",
                        "expected_impact": "How this will improve performance"
                    }},
                    {{
                        "suggestion": "Specific suggestion 5 for overall development",
                        "based_on": "Which feedback chunks support this suggestion",
                        "implementation": "Concrete steps to implement",
                        "expected_impact": "How this will improve performance"
                    }}
                ],
                "strengths": [
                    {{
                        "strength": "Specific strength 1 observed consistently",
                        "evidence": "Which chunks demonstrate this strength",
                        "frequency": "How often this strength appeared",
                        "impact": "How this strength supported communication"
                    }},
                    {{
                        "strength": "Specific strength 2 with pattern evidence", 
                        "evidence": "Which chunks demonstrate this strength",
                        "frequency": "How often this strength appeared",
                        "impact": "How this strength supported communication"
                    }},
                    {{
                        "strength": "Specific strength 3 showing development",
                        "evidence": "Which chunks demonstrate this strength",
                        "frequency": "How often this strength appeared",
                        "impact": "How this strength supported communication"
                    }}
                ],
                "practice_recommendations": [
                    {{
                        "activity": "Specific practice activity 1 targeting most common issue",
                        "priority": "High/Medium/Low priority based on frequency",
                        "frequency": "Recommended practice schedule",
                        "expected_outcome": "What improvement to expect"
                    }},
                    {{
                        "activity": "Specific practice activity 2 for pattern improvement",
                        "priority": "High/Medium/Low priority based on frequency",
                        "frequency": "Recommended practice schedule",
                        "expected_outcome": "What improvement to expect"
                    }},
                    {{
                        "activity": "Specific practice activity 3 for strength reinforcement", 
                        "priority": "High/Medium/Low priority based on frequency",
                        "frequency": "Recommended practice schedule",
                        "expected_outcome": "What improvement to expect"
                    }}
                ],
                "overall_scores": {{
                    "fluency": {fluency},
                    "pronunciation": {pronunciation},
                    "grammar": {grammar_score},
                    "emotion": "{emotion}",
                    "consistency_score": "X/10 (score for performance consistency across chunks)"
                }},
                "improvement_priority": {{
                    "area": "The single most important area to focus on based on pattern analysis",
                    "reason": "Why this area is critical (frequency, impact, etc.)",
                    "urgency": "High/Medium/Low based on analysis",
                    "first_steps": "Immediate actions to address this priority"
                }},
                "encouragement": {{
                    "progress_highlight": "Most significant improvement or consistent strength",
                    "pattern_achievement": "What the consistency patterns reveal about potential",
                    "motivational_message": "Personalized encouragement based on aggregate performance"
                }},
                "analysis_metadata": {{
                    "chunks_analyzed": "Number of feedback chunks processed",
                    "consistency_pattern": "Summary of performance consistency",
                    "key_insights": "Top 3 insights from pattern analysis"
                }}
            }}

            CRITICAL: Base all analysis on actual patterns found in the aggregated feedback chunks.
            Provide specific references to the chunk data where possible.
            Ensure all scores and feedback are justified by the aggregated evidence.
            """,
            input_variables=["final_feedback_result", "fluency", "pronunciation", "grammar_score", "emotion"]
        )

        parser = StrOutputParser()
        chain = prompt_template | model | parser

        feedback = await chain.ainvoke({
            "final_feedback_result": context_text,
            "fluency": fluency_score,
            "pronunciation": pronunciation_score,
            "grammar_score": grammar_score,
            "emotion": emotion_score
        })

        if "```json" in feedback:
            feedback = feedback.replace("```json", "")
        if "```" in feedback:
            feedback = feedback.replace("```", "")
        try:
            feedback_data = json.loads(feedback)
        except json.JSONDecodeError:
            feedback_data = {
                "raw_feedback": feedback,
                "overall_scores": {
                    "fluency": fluency_score,
                    "pronunciation": pronunciation_score,
                    "grammar": grammar_score,
                    "emotion": emotion_score
                }
            }

        logging.info(f"Overall scoring feedback generated successfully")

        if store_in_db and essay_id:
            await store_overall_scoring_in_db(essay_id, feedback_data)
        
        return feedback_data
        
    except ZeroDivisionError:
        logging.warning("No audio chunks were processed (count=0), skipping overall scoring")
        return None
    except Exception as e:
        logging.error(f"Error in overall scoring: {e}")
        import traceback
        traceback.print_exc()
        return None


async def store_overall_scoring_in_db(essay_id: str, scoring_data: Dict[str, Any]):

    try:
        essay_ref = db.collection("essays").document(essay_id)
        
        update_data = {
            "overall_scoring": scoring_data,
            "scoring_updated_at": datetime.now(),
            "status": "scoring_completed"
        }
        
        essay_ref.update(update_data)
        logging.info(f"Overall scoring stored for essay {essay_id}")
        logging.info("overall data stored in database ---------------------------->>>>>>>>>>>>>>")
        
    except Exception as e:
        logging.error(f"Failed to store overall scoring in database: {str(e)}")
        raise


def finalize_session(
    session_state, 
    username, 
    session_temp_dir, 
    topic, 
    essay_id, 
    chat_history, 
    final_feedback_result, 
    config,
):
    """
    Finalize session and trigger overall scoring in background
    """
    try:
        serializable_chat_history = []
        for message in chat_history:
            if isinstance(message, (AIMessage, SystemMessage, HumanMessage)):
                message_dict = {
                    'type': 'chat_message',
                    'content': message.content,
                    'message_type': 'system' if isinstance(message, SystemMessage) else
                                    'ai' if isinstance(message, AIMessage) else
                                    'human'
                }
                serializable_chat_history.append(message_dict)
            elif isinstance(message, dict):
                serializable_chat_history.append(message)
            else:
                serializable_chat_history.append(str(message))

        chunk_results_clean = []
        for chunk in session_state.get('chunk_results', []):
            clean_chunk = {
                'chunk_index': chunk.get('chunk_index'),
                'text': chunk.get('text'),
                'emotion': chunk.get('emotion'),
                'fluency': chunk.get('fluency'),
                'pronunciation': chunk.get('pronunciation'),
            }
            chunk_results_clean.append(clean_chunk)

        essay_ref = db.collection("essays").document(essay_id)
        update_data = {
            "chat_history": serializable_chat_history,
            "chunks": chunk_results_clean,
            "average_scores": topic.get_average_realtime_scores(),
            "updated_at": datetime.now(),
            "status": "completed",
            "transcript": " ".join(session_state.get('text_output', [])).strip()
        }
        essay_ref.update(update_data)

        asyncio.create_task(
            overall_scoring_by_speech_module_test(
                final_feedback_result=final_feedback_result,
                session_state=session_state,
                topic=topic,
                config=config,
                essay_id=essay_id,
                store_in_db=True
            )
        )
        
        logging.info(f"Successfully updated essay document {essay_id} and triggered scoring")
        return essay_id
        
    except Exception as e:
        logging.error(f"Failed to finalize essay document: {str(e)}", exc_info=True)
        raise


TEMP_DIR = os.path.abspath("temp_chunks")
os.makedirs(TEMP_DIR, exist_ok=True)

@app.websocket("/ws/assistant")
async def audio_ws_assistant(websocket: WebSocket):
    await websocket.accept()
    http_session = aiohttp.ClientSession()
    query_params = parse_qs(websocket.url.query)
    username = query_params.get("username", [None])[0]
    token = query_params.get("token", [None])[0]
    student_topic = query_params.get("topic", [None])[0]
    student_class = query_params.get("student_class", [None])[0]
    mood = query_params.get("mood", [None])[0]
    accent = query_params.get("accent", [None])[0]

    await validate_reset_token(token)

    final_feedback_result = []
    
    chat_history = []
    essay_id = await initialize_essay_document(
        username=username,
        student_topic=student_topic,
        student_class=student_class,
        mood=mood,
        accent=accent,
        chat_history=chat_history
    )

    logging.info(f"essay_id : {essay_id}")
    
    try:
        await websocket.send_json({"action": "essay_initialized", "essay_id": essay_id})
    except Exception as e:
        logging.error(f"Failed to send initial essay_id: {str(e)}")

    ai_response = await system_message(student_topic, mood, student_class, accent)
    await websocket.send_json({
                    "type": "ai_response",
                    "data": ai_response
                })
    chat_history.append(SystemMessage(content=ai_response))
    
    if not username or not token:
        await websocket.close(code=4001)
        logging.info("Username or token missing.")
        return

    logging.info(f"[WS] Authenticated connection from {username}")
    
    config = {
        'silence_threshold': 2,
        'min_utterance_length': 3,
        'min_speech_duration': 1,
        'silence_dBFS_threshold': -30,
        'processing_cooldown': 1.5,
        'max_repetitions': 2,
        'max_silence': 1.0,
        'chunk_duration': 0.5,
        'fluency_score': 0,
        'pronunciation_score': 0,
        'emotion_score': None,
        'grammar_score': 0,
        'count':0
    }

    session_state = {
        'silvero_model': True,
        'audio_buffer': AudioSegment.empty(),
        'speech_buffer': AudioSegment.empty(),
        'text_buffer': [],
        'last_speech_time': time.time(),
        'conversation_active': False,
        'processing_active': False,
        'assistant_speaking': False,
        'chunk_index': 0,
        'chunk_results': [],
        'text_output': [],
        'has_speech': False,
        'last_processing_time': 0,
        'consecutive_silence_chunks': 0,
        'active_speech_duration': 0.0
    }
    
    
    topic = Topic()
    session_temp_dir = tempfile.mkdtemp(prefix=f"{username}_", dir=TEMP_DIR)

    session_state["assistant_speaking"] = True
    response_audio = await topic.text_to_speech_assistant(ai_response, username, session_temp_dir)
    sleep_time = await send_audio_response(websocket, response_audio)
    await asyncio.sleep(sleep_time + 4)
    session_state["assistant_speaking"] = False

    logging.info(f"[Session] Temp dir created at {session_temp_dir}")
    
    final_output = os.path.join(session_temp_dir, f"{username}_output.wav")
    transcript_path = os.path.join(session_temp_dir, f"{username}_transcript.txt")
    
    audio_silence_handle = False
    silence_count = 0

    try:
        while True:
            message = await websocket.receive()

            if message["type"] == "websocket.disconnect":
                logging.info(f"[WS] {username} disconnected.")
                break

            if message["type"] == "websocket.receive" and "bytes" in message:
                if (session_state['processing_active'] or 
                    session_state['assistant_speaking'] or
                    (time.time() - session_state['last_processing_time'] < config['processing_cooldown'])):

                    current_cooldown = time.time() - session_state['last_processing_time']
                    if current_cooldown < config['processing_cooldown']:
                        logging.debug(f"In cooldown period: {current_cooldown:.2f}s remaining")
                    continue   

                current_time = time.time()
                new_chunk = AudioSegment(
                    data=message["bytes"],
                    sample_width=2,
                    frame_rate=16000,
                    channels=1
                )
                
                session_state['audio_buffer'] += new_chunk
                
                loudness = new_chunk.dBFS
                is_silent = loudness < config['silence_dBFS_threshold']

                if session_state["silvero_model"] == True:
                    if is_silent and audio_silence_handle == True:
                        silence_count += 1
                        session_state['consecutive_silence_chunks'] += 1
                        logging.info(f"🟡 Silent chunk detected (loudness: {loudness:.2f} dB)")
                        
                        if (session_state['has_speech'] and 
                            session_state['consecutive_silence_chunks'] * config['chunk_duration'] >= config['silence_threshold'] and
                            session_state['active_speech_duration'] >= config['min_speech_duration']) or silence_count == 3:
                            
                            await process_buffered_audio(
                                session_state, 
                                websocket, username, session_temp_dir, topic, 
                                config, student_topic, student_class, mood, accent, chat_history,http_session, final_feedback_result
                            )
                    else:
                        silence_count = 0
                        if current_time - session_state['last_processing_time'] < config['processing_cooldown']:
                            continue
                        session_state['consecutive_silence_chunks'] = 0
                        session_state['has_speech'] = True
                        session_state['speech_buffer'] += new_chunk
                        session_state['active_speech_duration'] += config['chunk_duration']
                        session_state['last_speech_time'] = current_time
                        logging.info(f"🔵 Speech chunk detected (loudness: {loudness:.2f} dB), active duration: {session_state['active_speech_duration']:.2f}s")
                        audio_silence_handle = True
                        
                        chunk_filename = os.path.join(session_temp_dir, f"chunk_temp_{session_state['chunk_index']}.wav")
                        new_chunk.export(chunk_filename, format="wav")
                        session_state['chunk_index'] += 1
                        
                        vad_result = await topic.silvero_vad(chunk_filename)
                        current_silence = vad_result.get("duration", 0.0)
                        
                        if current_silence > 0.3:
                            logging.info(f"VAD detected silence gap: {current_silence:.2f}s")
                            
                            if (current_silence >= config['silence_threshold'] and 
                                session_state['active_speech_duration'] >= config['min_speech_duration']):
                                
                                await process_buffered_audio(
                                    session_state, 
                                    websocket, username, session_temp_dir, topic, 
                                    config, student_topic, student_class, mood, accent, chat_history, http_session, final_feedback_result
                                )
    except WebSocketDisconnect:
        logging.warning(f"[WS] {username} disconnected unexpectedly")
    except Exception as e:
        logging.error(f"Unexpected error: {str(e)}", exc_info=True)
    finally:
        try:
            essay_id = finalize_session(session_state, username, session_temp_dir, topic, essay_id, chat_history, final_feedback_result, config)
        except Exception:
            pass
        with contextlib.suppress(Exception):
            await http_session.close()


def cleanup_temp_files(session_temp_dir, chunk_results):
    try:
        for chunk in chunk_results:
            try:
                if os.path.exists(chunk['file_path']):
                    os.remove(chunk['file_path'])
            except Exception as e:
                logging.warning(f"Failed to delete chunk file {chunk['file_path']}: {e}")

        try:
            if os.path.exists(session_temp_dir):
                if not os.listdir(session_temp_dir):
                    os.rmdir(session_temp_dir)
        except Exception as e:
            logging.warning(f"Failed to remove temp directory {session_temp_dir}: {e}")

    except Exception as e:
        logging.error(f"Error during temp file cleanup: {str(e)}")


async def process_buffered_audio(session_state, websocket, username, temp_dir, topic, config, student_topic, student_class, mood, accent, chat_history, http_session, final_feedback_result):
    if len(session_state['audio_buffer']) == 0:
        return
    
    is_greeting = False
        
    session_state['processing_active'] = True
    session_state["silvero_model"] = False
    
    try:
        session_state['consecutive_silence_chunks'] = 0
        session_state['has_speech'] = False
        session_state['active_speech_duration'] = 0.0
        session_state['last_speech_time'] = time.time()
        
        buffer_filename = os.path.join(temp_dir, f"buffered_{time.time()}.wav")
        session_state['audio_buffer'].export(buffer_filename, format="wav")
        
        transcribed_text = await topic.speech_to_text(buffer_filename, username)
        output = filter_to_english(transcribed_text)
        if not output:
            return
        
        await websocket.send_json({
                    "type": "transcribed_text",
                    "data": transcribed_text
                })
        logging.info(f"Transcribed: {transcribed_text}")
        session_state['text_buffer'].append(transcribed_text)
        session_state['text_output'].append(transcribed_text)

        test_list = list(transcribed_text.split())


        if len(test_list) > 10:
            is_greeting = True

        if is_greeting:
            asyncio.create_task(
                _safe_finalize_feedback(config, temp_dir, transcribed_text, websocket, session_state, username, buffer_filename, topic, http_session, final_feedback_result)
            )

        await process_user_utterance(
            transcribed_text,
            session_state, buffer_filename, websocket, 
            username, temp_dir, topic, student_topic, student_class, mood, accent, chat_history
        )

        session_state['audio_buffer'] = AudioSegment.empty()
        session_state['speech_buffer'] = AudioSegment.empty()
        session_state['text_buffer'] = []
        session_state['conversation_active'] = True
        
        session_state['last_processing_time'] = time.time()
        
    except Exception as e:
        logging.error(f"Error processing buffered audio: {str(e)}", exc_info=True)
    finally:
        session_state['processing_active'] = False
        await asyncio.sleep(0.5)
        session_state["silvero_model"] = True
        is_greeting = False
        logging.info(f"VAD model re-enabled, session state: {session_state['silvero_model']}")



async def scraping(topic: str) -> str:
    url = f"https://en.wikipedia.org/wiki/{topic}"
    api_endpoint = f"https://api.scrapingdog.com/scrape?api_key={scraping_api_key}&url={url}"

    text = ""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(api_endpoint) as response:
                if response.status == 200:
                    html = await response.text()
                    soup = BeautifulSoup(html, "html.parser")
                    for tag in soup(["script", "style", "noscript"]):
                        tag.decompose()
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    logging.warning(f"[SCRAPING] Failed with status {response.status}: {await response.text()}")
    except Exception as e:
        logging.exception(f"[SCRAPING ERROR] Failed to scrape data: {e}")

    return text



async def process_user_utterance(text, session_state, chunk_filename, websocket, 
                               username, session_temp_dir, topic,student_topic,student_class,mood,accent,chat_history):
    
    session_state['text_output'].append(text)
    session_state['chunk_index'] += 1
    scraped_data = await scraping(topic)

    try:
        prompt_template = PromptTemplate(template = """
            ROLE: You are a friendly, knowledgeable assistant and communicate like a friend. Your purpose is to:
            1. Answer questions conversationally
            2. Never reveal internal project details
            3. Keep responses under 200 words
            5. if by chance question is much more specific topic is needed and you does not have correct answer then used the function await scraping(your_specific_topic), which help in give updated answer.
            6. If you did not have updated data then always used the {scraped_data}, for update yourself, if by chance is question is for related to udpated data then also you able to send the request on the funtion await scraping(topic_you_want_to_select)                           
            Chat History:
            {chat_history}
            considering history as well before answering the question.

            Topic: {student_topic}
            Mood: {mood}
            Student Class: {student_class}
            Level: {level}
                                         
            - for reply anything consider always the Topic, Mood , student_class, Level and history.
            - And should be reply in the way and look like a friendly, knowledgeable teaching assistant.
            - Do not mention any technical details, model architecture, or team members.
            - Focus on providing helpful, concise answers.
            - If the question is not related to the topic, politely tell that quesiton is not related to the toic                     
                                                                
            RULES:
            - NEVER mention:
            * Model architecture/type
            * Team members/credentialsS
            * Code implementation
            * Technical specifications
            - Always redirect technical questions to general knowledge                                                                                                 

            USER QUESTION: {question}

            RESPONSE FORMAT:
            [Answer Concise 1-10 sentence response] 
            [if needed  Optional follow-up question to continue conversation]

            EXAMPLE:
            This exaple is for you never asked that. 
            User: What model are you using?
            I focus on helping with learning concepts rather than technical details. 
            Would you like me to explain how these systems generally work?

            Current response should be:
            """,input_variables=["scraped_data","question","student_topic","student_class","mood","accent","chat_history"]
            )

        model = ChatOllama(
                model="mistral",
                model_kwargs={"temperature": 0.8}
            )

        parser = StrOutputParser()

        chain = prompt_template | model | parser

        ai_response = await chain.ainvoke({
            "scraped_data": scraped_data,
            "student_topic": student_topic,
            "mood": mood,
            "student_class": student_class,
            "level": accent,
            "question": text,
            "chat_history": chat_history
        })
        chat_history.append(AIMessage(content=ai_response))
        await websocket.send_json({
                    "type": "ai_response",
                    "data": ai_response
                })
        logging.info(f"[AI Response]: {ai_response}")

        chunk_result = {
                    "chunk_index": session_state['chunk_index'],
                    "text": text,
                    "file_path": chunk_filename,
                    "chat_history":chat_history,
                }
        
        session_state["chunk_results"].append(chunk_result)

        response_audio = await topic.text_to_speech_assistant(ai_response, username, session_temp_dir)
        sleep_time = await send_audio_response(websocket, response_audio)
        await asyncio.sleep(sleep_time)

    except Exception as e:
        logging.error(f"QA Error: {str(e)}")
        await send_default_response(websocket, username, session_temp_dir, topic)


async def send_audio_response(websocket, audio_file):
    try:
        if not os.path.exists(audio_file):
            logging.error(f"Audio file not found: {audio_file}")
            return 0
            
        audio = AudioSegment.from_wav(audio_file)
        duration_ms = len(audio)
        
        if websocket.client_state == WebSocketState.DISCONNECTED:
            logging.warning("WebSocket disconnected before sending audio")
            return 0
            
        with open(audio_file, "rb") as f:
            await websocket.send_bytes(f.read())

        return duration_ms / 1000
        
    except Exception as e:
        logging.error(f"Failed to send audio response: {str(e)}")
        return 0


async def send_default_response(websocket, username, session_temp_dir, topic):
    try:
        default_response = "I didn't quite catch that. Could you please repeat?"
        response_audio = await topic.text_to_speech_assistant(default_response, username, session_temp_dir)
        await send_audio_response(websocket, response_audio)
    except Exception as e:
        logging.error(f"Failed to send default response: {str(e)}")

async def send_followup(websocket, username, session_temp_dir, topic):
    try:
        followup = "Is there anything else I can help you with?"
        response_audio = await topic.text_to_speech_assistant(followup, username, session_temp_dir)
        await send_audio_response(websocket, response_audio)
    except Exception as e:
        logging.error(f"Failed to send follow-up: {str(e)}")

